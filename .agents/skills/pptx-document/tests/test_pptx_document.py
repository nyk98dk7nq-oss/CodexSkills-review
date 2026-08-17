from __future__ import annotations

import json
from pathlib import Path
import subprocess
import sys
import tempfile
import unittest

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.util import Inches


SCRIPT = Path(__file__).resolve().parents[1] / "scripts" / "pptx_document.py"


def add_background_image(slide, image_path: Path) -> str:
    _, relationship_id = slide.part.get_or_add_image_part(str(image_path))
    background = parse_xml(
        f"""
        <p:bg {nsdecls('a', 'p', 'r')}>
          <p:bgPr>
            <a:blipFill>
              <a:blip r:embed="{relationship_id}"/>
              <a:stretch><a:fillRect/></a:stretch>
            </a:blipFill>
            <a:effectLst/>
          </p:bgPr>
        </p:bg>
        """
    )
    slide._element.cSld._remove_bg()
    slide._element.cSld._insert_bg(background)
    return relationship_id


class PptxDocumentCliTest(unittest.TestCase):
    def setUp(self) -> None:
        self.temporary = tempfile.TemporaryDirectory()
        self.root = Path(self.temporary.name)
        target_dir = self.root / "input" / "targets"
        target_dir.mkdir(parents=True)
        image_path = target_dir / "sample.png"
        Image.new("RGB", (32, 20), "green").save(image_path)
        self.source = target_dir / "design.pptx"

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        paragraph = text_box.text_frame.paragraphs[0]
        paragraph.add_run().text = "旧"
        paragraph.add_run().text = "名称の構成"
        self.text_shape_id = text_box.shape_id
        self.text_shape_name = text_box.name
        table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(4), Inches(1.5))
        table_shape.table.cell(0, 0).text = "項目"
        table_shape.table.cell(0, 1).text = "内容"
        table_shape.table.cell(1, 0).text = "方式"
        table_shape.table.cell(1, 1).text = "同期"
        picture_1 = slide.shapes.add_picture(
            str(image_path), Inches(6), Inches(1), width=Inches(0.5)
        )
        picture_2 = slide.shapes.add_picture(
            str(image_path), Inches(6), Inches(2), width=Inches(0.5)
        )
        self.picture_shape_ids = [picture_1.shape_id, picture_2.shape_id]
        slide.notes_slide.notes_text_frame.text = "レビュー用ノート"
        presentation.save(self.source)

    def tearDown(self) -> None:
        self.temporary.cleanup()

    def run_cli(self, *arguments: object) -> subprocess.CompletedProcess[str]:
        return subprocess.run(
            [sys.executable, str(SCRIPT), *(str(value) for value in arguments)],
            text=True,
            capture_output=True,
            check=False,
        )

    def make_symlink(
        self, link: Path, target: Path, *, target_is_directory: bool = False
    ) -> None:
        try:
            link.symlink_to(target, target_is_directory=target_is_directory)
        except (NotImplementedError, OSError) as exc:
            self.skipTest(f"シンボリックリンクを作成できません: {exc}")

    def test_to_markdown_extracts_shapes_table_notes_and_image(self) -> None:
        output = self.root / "work" / "markdown" / "design.md"
        images = self.root / "work" / "images"
        result = self.run_cli(
            "to-markdown",
            self.source,
            output,
            "--role",
            "target",
            "--repo-root",
            self.root,
            "--images-dir",
            images,
        )
        self.assertEqual(result.returncode, 0, result.stderr)
        text = output.read_text(encoding="utf-8")
        self.assertIn('source_path: "input/targets/design.pptx"', text)
        self.assertIn('document_role: "target"', text)
        self.assertIn('converter_skill: "pptx-document"', text)
        self.assertIn("## スライド 1", text)
        self.assertIn(f"### 図形 {self.text_shape_id}", text)
        self.assertIn("表 行2 列2", text)
        self.assertIn("レビュー用ノート", text)
        self.assertIn("![スライド1", text)
        self.assertEqual(len(list(images.glob("design-slide-001-shape-*.png"))), 2)

    def test_edit_replaces_across_runs_and_sets_shape_text(self) -> None:
        operations = self.root / "operations.json"
        operations.write_text(
            json.dumps(
                {
                    "operations": [
                        {"op": "replace_text", "old": "旧名称", "new": "新名称"},
                        {
                            "op": "set_shape_text",
                            "slide": 1,
                            "shape_name": self.text_shape_name,
                            "text": "確定した構成",
                        },
                    ]
                },
                ensure_ascii=False,
            ),
            encoding="utf-8",
        )
        output = self.root / "output" / "edited.pptx"
        result = self.run_cli("edit", self.source, output, "--operations", operations)
        self.assertEqual(result.returncode, 0, result.stderr)
        edited = Presentation(output)
        target = next(shape for shape in edited.slides[0].shapes if shape.shape_id == self.text_shape_id)
        self.assertEqual(target.text, "確定した構成")
        original = Presentation(self.source)
        original_target = next(
            shape for shape in original.slides[0].shapes if shape.shape_id == self.text_shape_id
        )
        self.assertEqual(original_target.text, "旧名称の構成")

    def test_edit_accepts_shape_id_for_backward_compatibility(self) -> None:
        operations = self.root / "operations.json"
        operations.write_text(
            json.dumps(
                [
                    {
                        "op": "set_shape_text",
                        "slide": 1,
                        "shape_id": self.text_shape_id,
                        "text": "IDによる更新",
                    }
                ],
                ensure_ascii=False,
            ),
            encoding="utf-8",
        )
        output = self.root / "output" / "edited-by-id.pptx"
        result = self.run_cli("edit", self.source, output, "--operations", operations)
        self.assertEqual(result.returncode, 0, result.stderr)
        edited = Presentation(output)
        target = next(
            shape for shape in edited.slides[0].shapes if shape.shape_id == self.text_shape_id
        )
        self.assertEqual(target.text, "IDによる更新")

    def test_edit_rejects_unknown_shape_name(self) -> None:
        operations = self.root / "operations.json"
        operations.write_text(
            json.dumps(
                [
                    {
                        "op": "set_shape_text",
                        "slide": 1,
                        "shape_name": "存在しない図形",
                        "text": "x",
                    }
                ],
                ensure_ascii=False,
            ),
            encoding="utf-8",
        )
        output = self.root / "output" / "edited.pptx"
        result = self.run_cli("edit", self.source, output, "--operations", operations)
        self.assertNotEqual(result.returncode, 0)
        self.assertIn("図形名が見つかりません", result.stderr)
        self.assertFalse(output.exists())

    def test_image_conflict_stops_before_writing_any_output(self) -> None:
        output = self.root / "work" / "markdown" / "design.md"
        images = self.source.parent
        first = images / (
            f"design-slide-001-shape-{self.picture_shape_ids[0]}.png"
        )
        conflict = images / (
            f"design-slide-001-shape-{self.picture_shape_ids[1]}.png"
        )
        conflict.write_bytes(b"sentinel")

        result = self.run_cli(
            "to-markdown",
            self.source,
            output,
            "--role",
            "target",
            "--repo-root",
            self.root,
            "--images-dir",
            images,
        )

        self.assertNotEqual(result.returncode, 0)
        self.assertIn("何も書き込みません", result.stderr)
        self.assertEqual(conflict.read_bytes(), b"sentinel")
        self.assertFalse(first.exists())
        self.assertFalse(output.exists())

    def test_broken_planned_image_symlink_stops_all_output(self) -> None:
        output = self.root / "work" / "markdown" / "design.md"
        images = self.root / "work" / "images"
        images.mkdir(parents=True)
        first = images / (
            f"design-slide-001-shape-{self.picture_shape_ids[0]}.png"
        )
        broken = images / (
            f"design-slide-001-shape-{self.picture_shape_ids[1]}.png"
        )
        self.make_symlink(broken, self.root / "outside" / "missing.png")

        result = self.run_cli(
            "to-markdown",
            self.source,
            output,
            "--role",
            "target",
            "--repo-root",
            self.root,
            "--images-dir",
            images,
        )

        self.assertNotEqual(result.returncode, 0)
        self.assertIn("シンボリックリンク", result.stderr)
        self.assertTrue(broken.is_symlink())
        self.assertFalse(first.exists())
        self.assertFalse(output.exists())

    def test_symlink_primary_output_parent_and_images_dir_are_rejected(self) -> None:
        work = self.root / "work"
        work.mkdir()
        broken_output = work / "broken.md"
        broken_target = self.root / "outside" / "missing.md"
        self.make_symlink(broken_output, broken_target)
        result = self.run_cli(
            "to-markdown",
            self.source,
            broken_output,
            "--role",
            "target",
            "--repo-root",
            self.root,
        )
        self.assertNotEqual(result.returncode, 0)
        self.assertTrue(broken_output.is_symlink())
        self.assertFalse(broken_target.exists())

        outside = self.root / "outside"
        outside.mkdir()
        output_parent_link = self.root / "edited-link"
        self.make_symlink(output_parent_link, outside, target_is_directory=True)
        operations = self.root / "empty-operations.json"
        operations.write_text("[]", encoding="utf-8")
        edited = output_parent_link / "edited.pptx"
        result = self.run_cli("edit", self.source, edited, "--operations", operations)
        self.assertNotEqual(result.returncode, 0)
        self.assertIn("シンボリックリンク", result.stderr)
        self.assertFalse((outside / "edited.pptx").exists())

        images_link = self.root / "images-link"
        self.make_symlink(images_link, outside, target_is_directory=True)
        markdown = work / "safe.md"
        result = self.run_cli(
            "to-markdown",
            self.source,
            markdown,
            "--role",
            "target",
            "--repo-root",
            self.root,
            "--images-dir",
            images_link,
        )
        self.assertNotEqual(result.returncode, 0)
        self.assertIn("シンボリックリンク", result.stderr)
        self.assertFalse(markdown.exists())
        self.assertEqual(list(outside.iterdir()), [])

    def test_commands_reject_existing_primary_output(self) -> None:
        markdown = self.root / "work" / "design.md"
        markdown.parent.mkdir(parents=True)
        markdown.write_text("sentinel", encoding="utf-8")
        images = self.root / "work" / "images"
        result = self.run_cli(
            "to-markdown",
            self.source,
            markdown,
            "--role",
            "target",
            "--repo-root",
            self.root,
            "--images-dir",
            images,
        )
        self.assertNotEqual(result.returncode, 0)
        self.assertEqual(markdown.read_text(encoding="utf-8"), "sentinel")
        self.assertFalse(images.exists())

        operations = self.root / "operations.json"
        operations.write_text(
            json.dumps(
                [
                    {
                        "op": "set_shape_text",
                        "slide": 1,
                        "shape_name": self.text_shape_name,
                        "text": "更新",
                    }
                ],
                ensure_ascii=False,
            ),
            encoding="utf-8",
        )
        edited = self.root / "output" / "edited.pptx"
        edited.parent.mkdir(parents=True)
        edited.write_bytes(b"sentinel")
        result = self.run_cli("edit", self.source, edited, "--operations", operations)
        self.assertNotEqual(result.returncode, 0)
        self.assertEqual(edited.read_bytes(), b"sentinel")

    def test_picture_placeholder_image_is_extracted(self) -> None:
        presentation = Presentation()
        picture_layout = next(
            (
                layout
                for layout in presentation.slide_layouts
                if any(
                    placeholder.placeholder_format.type == PP_PLACEHOLDER.PICTURE
                    for placeholder in layout.placeholders
                )
            ),
            None,
        )
        self.assertIsNotNone(picture_layout, "画像プレースホルダーつきレイアウトが必要です")
        slide = presentation.slides.add_slide(picture_layout)
        picture_placeholder = next(
            placeholder
            for placeholder in slide.placeholders
            if placeholder.placeholder_format.type == PP_PLACEHOLDER.PICTURE
        )
        source_image = self.source.parent / "sample.png"
        inserted = picture_placeholder.insert_picture(str(source_image))
        placeholder_source = self.source.parent / "placeholder.pptx"
        presentation.save(placeholder_source)

        output = self.root / "work" / "placeholder.md"
        images = self.root / "work" / "placeholder-images"
        result = self.run_cli(
            "to-markdown",
            placeholder_source,
            output,
            "--role",
            "target",
            "--repo-root",
            self.root,
            "--images-dir",
            images,
        )

        self.assertEqual(result.returncode, 0, result.stderr)
        text = output.read_text(encoding="utf-8")
        self.assertIn(f"図形 {inserted.shape_id}", text)
        self.assertIn("![スライド1", text)
        self.assertEqual(
            len(list(images.glob("placeholder-slide-001-shape-*.png"))), 1
        )

    def test_background_image_records_confirmation_marker_without_extraction(self) -> None:
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        relationship_id = add_background_image(
            slide, self.source.parent / "sample.png"
        )
        source = self.source.parent / "background.pptx"
        presentation.save(source)
        output = self.root / "work" / "background.md"
        images = self.root / "work" / "background-images"

        result = self.run_cli(
            "to-markdown",
            source,
            output,
            "--role",
            "target",
            "--repo-root",
            self.root,
            "--images-dir",
            images,
        )

        self.assertEqual(result.returncode, 0, result.stderr)
        text = output.read_text(encoding="utf-8")
        self.assertIn("背景画像は抽出対象外・要確認", text)
        self.assertIn("スライド 1/背景", text)
        self.assertIn("p:bg/p:bgPr/a:blipFill/a:blip", text)
        self.assertIn(f"embed={relationship_id}", text)
        self.assertIn("/ppt/slides/slide1.xml", text)
        self.assertFalse(images.exists())

    def test_solid_background_does_not_record_image_marker(self) -> None:
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x12, 0x34, 0x56)
        source = self.source.parent / "solid-background.pptx"
        presentation.save(source)
        output = self.root / "work" / "solid-background.md"

        result = self.run_cli(
            "to-markdown",
            source,
            output,
            "--role",
            "target",
            "--repo-root",
            self.root,
        )

        self.assertEqual(result.returncode, 0, result.stderr)
        text = output.read_text(encoding="utf-8")
        self.assertNotIn("背景画像は抽出対象外・要確認", text)
        self.assertNotIn("p:bg/p:bgPr/a:blipFill/a:blip", text)


if __name__ == "__main__":
    unittest.main()
