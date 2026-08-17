#!/usr/bin/env python3
"""PPTXを位置情報付きMarkdownへ変換し、安全な限定編集を行う。"""

from __future__ import annotations

import argparse
from datetime import datetime
import json
import os
from pathlib import Path
import sys
from typing import Any, Iterator


EMU_PER_INCH = 914400


def _require_runtime() -> None:
    if sys.version_info < (3, 12):
        raise RuntimeError(
            f"Python 3.12以上が必要です。現在: {sys.version.split()[0]}"
        )


def _require_pptx() -> Any:
    try:
        import pptx
    except ImportError as exc:
        raise RuntimeError(
            "必要なライブラリ python-pptx を読み込めません。READMEのインストール手順を確認してください。"
        ) from exc
    return pptx


def _validate_input(path: Path) -> None:
    if not path.is_file():
        raise FileNotFoundError(f"入力ファイルが見つかりません: {path}")
    if path.suffix.lower() != ".pptx":
        raise ValueError(f"PPTXファイルを指定してください: {path}")


def _ensure_distinct(input_path: Path, output_path: Path) -> None:
    if input_path.resolve() == output_path.resolve():
        raise ValueError("入力ファイルは上書きできません。別の出力パスを指定してください。")
    _validate_new_output_path(output_path, "出力先")


def _reject_symlink_components(path: Path, label: str) -> None:
    """出力パス自体と既存祖先のsymbolic linkを拒否する。"""
    absolute = path if path.is_absolute() else Path.cwd() / path
    for candidate in (absolute, *absolute.parents):
        if _is_link_like(candidate):
            raise ValueError(
                f"{label}にシンボリックリンクを使用できません: {candidate}"
            )


def _is_link_like(path: Path) -> bool:
    if path.is_symlink():
        return True
    is_junction = getattr(path, "is_junction", None)
    if is_junction is None:
        return False
    try:
        return bool(is_junction())
    except OSError:
        return False


def _validate_new_output_path(path: Path, label: str) -> None:
    _reject_symlink_components(path, label)
    if os.path.lexists(path):
        raise FileExistsError(
            f"{label}は既に存在します。上書きせず、何も書き込みません: {path}"
        )


def _validate_output_directory(path: Path, label: str) -> None:
    _reject_symlink_components(path, label)
    if os.path.lexists(path) and not path.is_dir():
        raise NotADirectoryError(f"{label}はディレクトリではありません: {path}")


def _repo_relative(path: Path, repo_root: Path) -> str:
    try:
        return path.resolve().relative_to(repo_root.resolve()).as_posix()
    except ValueError as exc:
        raise ValueError(
            f"入力ファイルはリポジトリルート配下に置いてください: {path}"
        ) from exc


def _yaml_string(value: str) -> str:
    return json.dumps(value, ensure_ascii=False)


def _md(value: Any) -> str:
    text = "" if value is None else str(value)
    return text.replace("\r\n", "\n").replace("\r", "\n").replace("|", "\\|").replace("\n", "<br>")


def _walk_shapes(shapes: Any, prefix: str = "") -> Iterator[tuple[Any, str]]:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for position, shape in enumerate(shapes, start=1):
        path = f"{prefix}{position}" if not prefix else f"{prefix}.{position}"
        yield shape, path
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _walk_shapes(shape.shapes, path)


def _shape_bbox(shape: Any) -> str:
    return (
        f"left={shape.left / EMU_PER_INCH:.2f}in, "
        f"top={shape.top / EMU_PER_INCH:.2f}in, "
        f"width={shape.width / EMU_PER_INCH:.2f}in, "
        f"height={shape.height / EMU_PER_INCH:.2f}in"
    )


def _text_frame_lines(text_frame: Any, location: str) -> list[str]:
    lines: list[str] = []
    for paragraph_index, paragraph in enumerate(text_frame.paragraphs, start=1):
        if not paragraph.text:
            continue
        lines.extend(
            [
                f"- 位置: `{location}/段落 {paragraph_index}`; 箇条書きレベル: {paragraph.level}",
                f"  {_md(paragraph.text)}",
            ]
        )
    return lines


def _background_owners(slide: Any) -> Iterator[tuple[str, Any]]:
    """スライドからレイアウト、マスターの順に背景候補を返す。"""
    yield "スライド", slide
    slide_layout = getattr(slide, "slide_layout", None)
    if slide_layout is None:
        return
    yield "スライドレイアウト", slide_layout
    slide_master = getattr(slide_layout, "slide_master", None)
    if slide_master is not None:
        yield "スライドマスター", slide_master


def _background_image_evidence(slide: Any) -> list[dict[str, str]]:
    """
    実効背景の `p:bg/p:bgPr/a:blipFill/a:blip` だけを画像として検出する。

    下位の明示背景があればそこで継承を止め、単色・グラデーション背景に
    対してマスター画像を誤検出しない。
    """
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    from pptx.oxml.ns import qn

    for owner_name, owner in _background_owners(slide):
        common_slide_data = getattr(getattr(owner, "_element", None), "cSld", None)
        background = getattr(common_slide_data, "bg", None)
        if background is None:
            continue

        part = getattr(owner, "part", None)
        if part is None:
            return []
        evidence: list[dict[str, str]] = []
        for blip in background.xpath("./p:bgPr/a:blipFill/a:blip"):
            for relationship_kind, attribute in (
                ("embed", qn("r:embed")),
                ("link", qn("r:link")),
            ):
                relationship_id = blip.get(attribute)
                if not relationship_id:
                    continue
                try:
                    relationship = part.rels[relationship_id]
                except KeyError:
                    evidence.append(
                        {
                            "owner": owner_name,
                            "part": str(getattr(part, "partname", "不明")),
                            "relationship_id": relationship_id,
                            "relationship_kind": relationship_kind,
                            "target": "relationship未解決",
                        }
                    )
                    continue
                if relationship.reltype != RT.IMAGE:
                    continue
                evidence.append(
                    {
                        "owner": owner_name,
                        "part": str(getattr(part, "partname", "不明")),
                        "relationship_id": relationship_id,
                        "relationship_kind": relationship_kind,
                        "target": str(getattr(relationship, "target_ref", "不明")),
                    }
                )
        # 明示された背景は、画像でなくても上位背景を上書きする。
        return evidence
    return []


def _background_markdown(slide: Any, slide_index: int) -> list[str]:
    evidence = _background_image_evidence(slide)
    if not evidence:
        return []
    lines = [
        "### 背景画像（要確認）",
        "",
        f"- 位置: `スライド {slide_index}/背景`",
        "- 判定: **背景画像は抽出対象外・要確認**",
    ]
    for item in evidence:
        lines.append(
            "- 根拠: "
            "`p:bg/p:bgPr/a:blipFill/a:blip` に"
            f" `{item['relationship_kind']}={item['relationship_id']}` を検出；"
            f" 由来 `{item['owner']}`；パート `{item['part']}`；"
            f" 参照先 `{item['target']}`"
        )
    lines.append("")
    return lines


def _shape_image(shape: Any) -> Any | None:
    """通常画像と画像プレースホルダーから画像を取得する。"""
    try:
        return shape.image
    except (AttributeError, ValueError):
        return None


def _plan_images(
    presentation: Any,
    images_dir: Path | None,
    source_stem: str,
) -> tuple[dict[tuple[int, int], dict[str, Any]], list[dict[str, Any]]]:
    """画像出力を一括計画し、既存パスや同一パスの衝突を書き込み前に拒否する。"""
    by_shape: dict[tuple[int, int], dict[str, Any]] = {}
    plans: list[dict[str, Any]] = []
    seen_paths: set[str] = set()
    conflicts: list[Path] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape, _ in _walk_shapes(slide.shapes):
            image = _shape_image(shape)
            if image is None:
                continue
            key = (slide_index, shape.shape_id)
            plan: dict[str, Any] = {
                "slide": slide_index,
                "shape_id": shape.shape_id,
                "path": None,
                "blob": None,
            }
            if images_dir is not None:
                extension = image.ext or "bin"
                image_path = images_dir / (
                    f"{source_stem}-slide-{slide_index:03d}-shape-{shape.shape_id}.{extension}"
                )
                _validate_new_output_path(image_path, "抽出画像の出力先")
                normalized = str(image_path.resolve()).casefold()
                if normalized in seen_paths:
                    conflicts.append(image_path)
                seen_paths.add(normalized)
                plan["path"] = image_path
                plan["blob"] = image.blob
                plans.append(plan)
            by_shape[key] = plan

    if conflicts:
        paths = ", ".join(str(path) for path in conflicts)
        raise FileExistsError(
            "画像出力先が既に存在するか、出力予定パスが重複しています。"
            f"何も書き込みません: {paths}"
        )
    return by_shape, plans


def _picture_markdown(
    plan: dict[str, Any],
    output_path: Path,
) -> str:
    slide_index = plan["slide"]
    shape_id = plan["shape_id"]
    image_path = plan["path"]
    if image_path is None:
        return f"[埋め込み画像: 図形ID {shape_id}]"
    relative_link = Path(os.path.relpath(image_path, output_path.parent)).as_posix()
    return f"![スライド{slide_index} 図形{shape_id}]({relative_link})"


def to_markdown(
    input_path: Path,
    output_path: Path,
    role: str,
    repo_root: Path,
    images_dir: Path | None,
) -> None:
    pptx = _require_pptx()

    _validate_input(input_path)
    _ensure_distinct(input_path, output_path)
    source_path = _repo_relative(input_path, repo_root)
    if images_dir is not None:
        _validate_output_directory(images_dir, "画像出力ディレクトリ")
    presentation = pptx.Presentation(input_path)
    image_by_shape, image_plans = _plan_images(
        presentation, images_dir, input_path.stem
    )
    lines = [
        "---",
        f"source_path: {_yaml_string(source_path)}",
        f"source_name: {_yaml_string(input_path.name)}",
        'source_format: "pptx"',
        f"document_role: {_yaml_string(role)}",
        f"converted_at: {_yaml_string(datetime.now().astimezone().isoformat(timespec='seconds'))}",
        'converter_skill: "pptx-document"',
        "---",
        "",
        f"# {_md(input_path.name)}",
        "",
    ]

    for slide_index, slide in enumerate(presentation.slides, start=1):
        lines.extend([f"## スライド {slide_index}", ""])
        lines.extend(_background_markdown(slide, slide_index))
        for shape, tree_path in _walk_shapes(slide.shapes):
            location = f"スライド {slide_index}/図形ID {shape.shape_id}"
            lines.extend(
                [
                    f"### 図形 {shape.shape_id}",
                    "",
                    f"- 位置: `{location}`",
                    f"- 図形ツリー位置: `{tree_path}`",
                    f"- 図形名: {_md(shape.name)}",
                    f"- 座標: `{_shape_bbox(shape)}`",
                    f"- 種別: `{shape.shape_type}`",
                    "",
                ]
            )
            image_plan = image_by_shape.get((slide_index, shape.shape_id))
            if image_plan is not None:
                lines.extend(
                    [
                        _picture_markdown(image_plan, output_path),
                        "",
                    ]
                )
            if getattr(shape, "has_text_frame", False):
                text_lines = _text_frame_lines(shape.text_frame, location)
                lines.extend(text_lines if text_lines else ["（テキストなし）"])
                lines.append("")
            if getattr(shape, "has_table", False):
                table = shape.table
                width = len(table.columns)
                lines.extend(
                    [
                        f"#### 表（{location}）",
                        "",
                        "| " + " | ".join(f"列 {index}" for index in range(1, width + 1)) + " |",
                        "| " + " | ".join("---" for _ in range(width)) + " |",
                    ]
                )
                for row_index, row in enumerate(table.rows, start=1):
                    values = [
                        f"[{location}/表 行{row_index} 列{column_index}] {_md(cell.text)}"
                        for column_index, cell in enumerate(row.cells, start=1)
                    ]
                    lines.append("| " + " | ".join(values) + " |")
                lines.append("")

        if slide.has_notes_slide:
            notes_text_frame = slide.notes_slide.notes_text_frame
            if notes_text_frame is not None and notes_text_frame.text.strip():
                lines.extend(["### ノート", "", f"- 位置: `スライド {slide_index}/ノート`", ""])
                lines.extend(_text_frame_lines(notes_text_frame, f"スライド {slide_index}/ノート"))
                lines.append("")

    for plan in image_plans:
        image_path = plan["path"]
        image_path.parent.mkdir(parents=True, exist_ok=True)
        with image_path.open("xb") as image_file:
            image_file.write(plan["blob"])
    output_path.parent.mkdir(parents=True, exist_ok=True)
    with output_path.open("x", encoding="utf-8") as output_file:
        output_file.write("\n".join(lines).rstrip() + "\n")


def _load_operations(path: Path) -> list[dict[str, Any]]:
    if not path.is_file():
        raise FileNotFoundError(f"操作JSONが見つかりません: {path}")
    try:
        payload = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as exc:
        raise ValueError(f"操作JSONを読み取れません: {path}: {exc}") from exc
    operations = payload.get("operations") if isinstance(payload, dict) else payload
    if not isinstance(operations, list) or not all(isinstance(item, dict) for item in operations):
        raise ValueError("操作JSONはオブジェクトの配列、または operations 配列を持つオブジェクトにしてください。")
    return operations


def _shape_by_id(shapes: Any, shape_id: int) -> Any | None:
    for shape, _ in _walk_shapes(shapes):
        if shape.shape_id == shape_id:
            return shape
    return None


def _shape_by_selector(
    shapes: Any,
    shape_name: str | None,
    shape_id: int | None,
) -> Any:
    if shape_name is not None:
        if not isinstance(shape_name, str) or not shape_name:
            raise ValueError("shape_name は空でない文字列にしてください。")
        matches = [shape for shape, _ in _walk_shapes(shapes) if shape.name == shape_name]
        if shape_id is not None:
            if not isinstance(shape_id, int):
                raise ValueError("shape_id は整数にしてください。")
            matches = [shape for shape in matches if shape.shape_id == shape_id]
        if not matches:
            qualifier = f" / 図形ID {shape_id}" if shape_id is not None else ""
            raise ValueError(f"図形名が見つかりません: {shape_name}{qualifier}")
        if len(matches) > 1:
            raise ValueError(
                f"図形名が重複しています: {shape_name}。shape_id も指定してください。"
            )
        return matches[0]

    if shape_id is None:
        raise ValueError(
            "shape_name、または後方互換用の shape_id を指定してください。"
        )
    if not isinstance(shape_id, int):
        raise ValueError("shape_id は整数にしてください。")
    shape = _shape_by_id(shapes, shape_id)
    if shape is None:
        raise ValueError(f"図形IDが見つかりません: {shape_id}")
    return shape


def _all_text_frames(presentation: Any) -> Iterator[Any]:
    for slide in presentation.slides:
        for shape, _ in _walk_shapes(slide.shapes):
            if getattr(shape, "has_text_frame", False):
                yield shape.text_frame
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        yield cell.text_frame
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            yield slide.notes_slide.notes_text_frame


def _replace_paragraph(paragraph: Any, old: str, new: str, limit: int | None) -> int:
    current = paragraph.text
    if not old or old not in current or limit == 0:
        return 0
    occurrences = current.count(old)
    replacement_count = occurrences if limit is None else min(occurrences, limit)
    replaced = current.replace(old, new, replacement_count)
    if paragraph.runs:
        paragraph.runs[0].text = replaced
        for run in paragraph.runs[1:]:
            run.text = ""
    else:
        paragraph.text = replaced
    return replacement_count


def edit(input_path: Path, output_path: Path, operations_path: Path) -> None:
    pptx = _require_pptx()
    _validate_input(input_path)
    _ensure_distinct(input_path, output_path)
    presentation = pptx.Presentation(input_path)

    for index, operation in enumerate(_load_operations(operations_path), start=1):
        operation_name = operation.get("op")
        try:
            if operation_name == "replace_text":
                old = operation["old"]
                new = operation["new"]
                if not isinstance(old, str) or not old:
                    raise ValueError("old は空でない文字列にしてください。")
                if not isinstance(new, str):
                    raise ValueError("new は文字列にしてください。")
                count = operation.get("count")
                if count is not None and (not isinstance(count, int) or count < 1):
                    raise ValueError("count は1以上の整数にしてください。")
                remaining = count
                replacements = 0
                for text_frame in _all_text_frames(presentation):
                    for paragraph in text_frame.paragraphs:
                        changed = _replace_paragraph(paragraph, old, new, remaining)
                        replacements += changed
                        if remaining is not None:
                            remaining -= changed
                            if remaining == 0:
                                break
                    if remaining == 0:
                        break
                if replacements == 0:
                    raise ValueError(f"置換対象が見つかりません: {old}")
            elif operation_name == "set_shape_text":
                slide_index = operation["slide"]
                shape_name = operation.get("shape_name")
                shape_id = operation.get("shape_id")
                text = operation["text"]
                if not isinstance(slide_index, int) or slide_index < 1:
                    raise ValueError("slide は1以上の整数にしてください。")
                if not isinstance(text, str):
                    raise ValueError("text は文字列にしてください。")
                slide = presentation.slides[slide_index - 1]
                shape = _shape_by_selector(slide.shapes, shape_name, shape_id)
                if not getattr(shape, "has_text_frame", False):
                    raise ValueError(
                        f"図形 {shape.name} (ID {shape.shape_id}) はテキスト図形ではありません。"
                    )
                shape.text_frame.text = text
            else:
                raise ValueError(f"許可されていない操作です: {operation_name}")
        except (IndexError, KeyError, TypeError, ValueError) as exc:
            raise ValueError(f"操作 {index} が不正です: {exc}") from exc

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description=__doc__)
    subparsers = parser.add_subparsers(dest="command", required=True)
    markdown_parser = subparsers.add_parser("to-markdown", help="PPTXをMarkdownへ変換する")
    markdown_parser.add_argument("input", type=Path)
    markdown_parser.add_argument("output", type=Path)
    markdown_parser.add_argument("--role", required=True, choices=("checklist", "reference", "target"))
    markdown_parser.add_argument("--repo-root", required=True, type=Path)
    markdown_parser.add_argument("--images-dir", type=Path)

    edit_parser = subparsers.add_parser("edit", help="許可した操作だけでPPTXを編集する")
    edit_parser.add_argument("input", type=Path)
    edit_parser.add_argument("output", type=Path)
    edit_parser.add_argument("--operations", required=True, type=Path)
    return parser


def main() -> int:
    try:
        _require_runtime()
        arguments = build_parser().parse_args()
        if arguments.command == "to-markdown":
            to_markdown(
                arguments.input,
                arguments.output,
                arguments.role,
                arguments.repo_root,
                arguments.images_dir,
            )
        else:
            edit(arguments.input, arguments.output, arguments.operations)
        return 0
    except Exception as exc:
        print(f"エラー: {exc}", file=sys.stderr)
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
