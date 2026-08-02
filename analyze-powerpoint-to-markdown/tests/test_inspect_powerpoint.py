from __future__ import annotations

import contextlib
import hashlib
import importlib.util
import io
import json
import stat
import subprocess
import sys
import tempfile
import unittest
import zipfile
import xml.etree.ElementTree as ElementTree
from pathlib import Path
from unittest import mock

from PIL import Image
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.util import Inches


SCRIPT = Path(__file__).resolve().parents[1] / "scripts" / "inspect_powerpoint.py"
REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
CONTENT_TYPES_NS = "http://schemas.openxmlformats.org/package/2006/content-types"
DRAWING_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
OFFICE_REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
SVG_NS = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"
IMAGE_REL_TYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
)


def local_name(tag: str) -> str:
    return tag.rsplit("}", 1)[-1] if "}" in tag else tag


def c_nv_pr(shape):
    for element in shape.element.iter():
        if local_name(element.tag) == "cNvPr":
            return element
    raise AssertionError("cNvPr not found")


def rewrite_zip(
    source: Path,
    target: Path,
    *,
    replacements: dict[str, bytes] | None = None,
    additions: list[tuple[zipfile.ZipInfo | str, bytes]] | None = None,
) -> None:
    replacements = replacements or {}
    additions = additions or []
    with zipfile.ZipFile(source, "r") as input_archive, zipfile.ZipFile(
        target, "w", zipfile.ZIP_DEFLATED
    ) as output_archive:
        for info in input_archive.infolist():
            data = replacements.get(info.filename, input_archive.read(info.filename))
            output_archive.writestr(info, data)
        for name_or_info, data in additions:
            output_archive.writestr(name_or_info, data)


def iter_shapes(shapes):
    for shape in shapes:
        yield shape
        yield from iter_shapes(shape.get("children", []))


def retarget_first_picture(
    source: Path,
    target: Path,
    *,
    filename: str | None = None,
    content_type: str | None = None,
    data: bytes | None = None,
    external_target: str | None = None,
) -> None:
    relationship_part = "ppt/slides/_rels/slide1.xml.rels"
    with zipfile.ZipFile(source, "r") as archive:
        relationships = archive.read(relationship_part)
        content_types = archive.read("[Content_Types].xml")

    relationships_root = ElementTree.fromstring(relationships)
    image_relationship = next(
        element
        for element in relationships_root
        if element.attrib.get("Type") == IMAGE_REL_TYPE
    )
    replacements = {
        relationship_part: b"",
    }
    additions: list[tuple[zipfile.ZipInfo | str, bytes]] = []
    if external_target is not None:
        image_relationship.set("Target", external_target)
        image_relationship.set("TargetMode", "External")
    else:
        assert filename is not None and content_type is not None and data is not None
        image_relationship.set("Target", f"../media/{filename}")
        image_relationship.attrib.pop("TargetMode", None)
        content_types_root = ElementTree.fromstring(content_types)
        extension = Path(filename).suffix.lstrip(".")
        ElementTree.SubElement(
            content_types_root,
            f"{{{CONTENT_TYPES_NS}}}Default",
            {"Extension": extension, "ContentType": content_type},
        )
        replacements["[Content_Types].xml"] = ElementTree.tostring(
            content_types_root, encoding="utf-8", xml_declaration=True
        )
        additions.append((f"ppt/media/{filename}", data))
    replacements[relationship_part] = ElementTree.tostring(
        relationships_root, encoding="utf-8", xml_declaration=True
    )
    rewrite_zip(source, target, replacements=replacements, additions=additions)


class InspectPowerPointTests(unittest.TestCase):
    def setUp(self) -> None:
        self.temporary = tempfile.TemporaryDirectory()
        self.root = Path(self.temporary.name)
        self.presentation = self.root / "sample.pptx"
        self._create_sample(self.presentation)

    def tearDown(self) -> None:
        self.temporary.cleanup()

    def _create_sample(self, path: Path) -> None:
        image_path = self.root / "visible.png"
        Image.new("RGB", (32, 20), (30, 90, 180)).save(image_path)

        presentation = Presentation()
        presentation.core_properties.author = "SECRET-DOCUMENT-AUTHOR"
        presentation.core_properties.title = "SECRET-DOCUMENT-TITLE"

        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "VISIBLE-TITLE"
        body = slide.placeholders[1].text_frame
        body.clear()
        first = body.paragraphs[0]
        first.text = "VISIBLE-BODY"
        second = body.add_paragraph()
        second.text = "VISIBLE-CHILD"
        second.level = 1
        run = second.add_run()
        run.text = " VISIBLE-LINK"
        run.hyperlink.address = "https://example.invalid/reference"

        table_shape = slide.shapes.add_table(2, 2, Inches(0.5), Inches(3.2), Inches(4.0), Inches(1.2))
        table = table_shape.table
        table.cell(0, 0).text = "VISIBLE-TABLE-HEADER"
        table.cell(0, 0).merge(table.cell(0, 1))
        table.cell(1, 0).text = "VISIBLE-R2C1"
        table.cell(1, 1).text = "VISIBLE-R2C2"

        chart_data = CategoryChartData()
        chart_data.categories = ["VISIBLE-Q1", "VISIBLE-Q2"]
        chart_data.add_series("VISIBLE-SERIES", (10, 20))
        slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(5.0),
            Inches(0.8),
            Inches(4.0),
            Inches(2.5),
            chart_data,
        )

        picture = slide.shapes.add_picture(
            str(image_path), Inches(5.0), Inches(3.6), width=Inches(1.6)
        )
        c_nv_pr(picture).set("descr", "VISIBLE-IMAGE-ALT")

        hidden_shape = slide.shapes.add_textbox(Inches(0.5), Inches(5.1), Inches(3.0), Inches(0.5))
        hidden_shape.text = "SECRET-HIDDEN-SHAPE"
        c_nv_pr(hidden_shape).set("hidden", "1")

        hidden_group = slide.shapes.add_group_shape()
        hidden_group_child = hidden_group.shapes.add_textbox(
            Inches(0.5), Inches(5.7), Inches(3.0), Inches(0.4)
        )
        hidden_group_child.text = "SECRET-HIDDEN-GROUP-CHILD"
        c_nv_pr(hidden_group).set("hidden", "1")

        off_slide = slide.shapes.add_textbox(
            presentation.slide_width + Inches(1.0),
            Inches(1.0),
            Inches(2.0),
            Inches(0.5),
        )
        off_slide.text = "SECRET-OFF-SLIDE"

        notes = slide.notes_slide.notes_text_frame
        notes.text = "SECRET-SPEAKER-NOTE"

        hidden_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        hidden_slide.shapes.title.text = "SECRET-HIDDEN-SLIDE"
        hidden_slide.placeholders[1].text = "SECRET-HIDDEN-SLIDE-BODY"
        hidden_slide.element.set("show", "0")

        presentation.save(path)

    def run_inspector(self, path: Path | None = None, *arguments: str):
        source = path or self.presentation
        output = self.root / f"output-{len(list(self.root.glob('output-*.json')))}.json"
        command = [sys.executable, str(SCRIPT), str(source), "--output", str(output), *arguments]
        completed = subprocess.run(command, capture_output=True, text=True, check=False)
        data = json.loads(output.read_text(encoding="utf-8")) if output.exists() else None
        raw = output.read_text(encoding="utf-8") if output.exists() else ""
        return completed, data, raw, output

    def test_default_excludes_sensitive_content_and_extracts_visible_structure(self) -> None:
        before = hashlib.sha256(self.presentation.read_bytes()).hexdigest()
        completed, data, raw, _ = self.run_inspector()
        after = hashlib.sha256(self.presentation.read_bytes()).hexdigest()

        self.assertEqual(completed.returncode, 0, completed.stderr)
        self.assertEqual(before, after)
        self.assertIn("VISIBLE-TITLE", raw)
        self.assertIn("VISIBLE-TABLE-HEADER", raw)
        self.assertIn("VISIBLE-SERIES", raw)
        self.assertIn("VISIBLE-IMAGE-ALT", raw)
        self.assertIn("https://example.invalid/reference", raw)
        for secret in (
            "SECRET-HIDDEN-SHAPE",
            "SECRET-HIDDEN-GROUP-CHILD",
            "SECRET-OFF-SLIDE",
            "SECRET-SPEAKER-NOTE",
            "SECRET-HIDDEN-SLIDE",
            "SECRET-DOCUMENT-AUTHOR",
            "SECRET-DOCUMENT-TITLE",
        ):
            self.assertNotIn(secret, raw)

        assert data is not None
        self.assertEqual(data["presentation"]["slide_count"], 2)
        self.assertEqual(data["presentation"]["included_slide_count"], 1)
        self.assertEqual(data["presentation"]["hidden_slide_count"], 1)
        self.assertTrue(data["slides"][0]["has_notes"])
        self.assertFalse(data["slides"][0]["notes_included"])
        self.assertIsNone(data["slides"][0]["notes"])
        self.assertFalse(data["slides"][1]["included"])
        self.assertEqual(data["slides"][1]["exclusion_reason"], "hidden_slide")

        shapes = list(iter_shapes(data["slides"][0]["shapes"]))
        self.assertTrue(any(shape.get("table") for shape in shapes))
        self.assertTrue(any(shape.get("chart") for shape in shapes))
        self.assertTrue(any(shape.get("image") for shape in shapes))
        self.assertTrue(any(shape.get("exclusion_reason") == "hidden_shape" for shape in shapes))
        self.assertTrue(any(shape.get("exclusion_reason") == "fully_off_slide" for shape in shapes))

    def test_hidden_relationship_type_is_bucketed_without_string_leak(self) -> None:
        relationship_part = "ppt/slides/_rels/slide2.xml.rels"
        with zipfile.ZipFile(self.presentation, "r") as archive:
            relationships = archive.read(relationship_part)
        root = ElementTree.fromstring(relationships)
        ElementTree.SubElement(
            root,
            f"{{{REL_NS}}}Relationship",
            {
                "Id": "rIdHiddenSecretType",
                "Type": "http://example.invalid/relationships/SECRET-HIDDEN-RELTYPE",
                "Target": "../presentation.xml",
            },
        )
        modified = self.root / "hidden-relationship-type.pptx"
        rewrite_zip(
            self.presentation,
            modified,
            replacements={
                relationship_part: ElementTree.tostring(
                    root, encoding="utf-8", xml_declaration=True
                )
            },
        )

        completed, data, raw, _ = self.run_inspector(modified)
        self.assertEqual(completed.returncode, 0, completed.stderr)
        self.assertNotIn("SECRET-HIDDEN-RELTYPE", raw)
        assert data is not None
        self.assertGreaterEqual(
            data["presentation"]["archive"]["relationship_types"].get("other", 0),
            1,
        )

    def test_sensitive_options_are_independent(self) -> None:
        completed, _, raw, _ = self.run_inspector(None, "--include-notes")
        self.assertEqual(completed.returncode, 0, completed.stderr)
        self.assertIn("SECRET-SPEAKER-NOTE", raw)
        self.assertNotIn("SECRET-HIDDEN-SHAPE", raw)
        self.assertNotIn("SECRET-OFF-SLIDE", raw)
        self.assertNotIn("SECRET-HIDDEN-SLIDE", raw)
        self.assertNotIn("SECRET-DOCUMENT-AUTHOR", raw)

        completed, _, raw, _ = self.run_inspector(None, "--include-hidden-shapes")
        self.assertEqual(completed.returncode, 0, completed.stderr)
        self.assertIn("SECRET-HIDDEN-SHAPE", raw)
        self.assertIn("SECRET-HIDDEN-GROUP-CHILD", raw)
        self.assertNotIn("SECRET-SPEAKER-NOTE", raw)
        self.assertNotIn("SECRET-OFF-SLIDE", raw)
        self.assertNotIn("SECRET-HIDDEN-SLIDE", raw)

        completed, _, raw, _ = self.run_inspector(None, "--include-off-slide-shapes")
        self.assertEqual(completed.returncode, 0, completed.stderr)
        self.assertIn("SECRET-OFF-SLIDE", raw)
        self.assertNotIn("SECRET-HIDDEN-SHAPE", raw)
        self.assertNotIn("SECRET-SPEAKER-NOTE", raw)

        completed, _, raw, _ = self.run_inspector(None, "--include-document-properties")
        self.assertEqual(completed.returncode, 0, completed.stderr)
        self.assertIn("SECRET-DOCUMENT-AUTHOR", raw)
        self.assertNotIn("SECRET-HIDDEN-SHAPE", raw)
        self.assertNotIn("SECRET-SPEAKER-NOTE", raw)

    def test_hidden_slide_requires_explicit_selection_or_flag(self) -> None:
        completed, data, raw, _ = self.run_inspector(None, "--slide", "2")
        self.assertEqual(completed.returncode, 0, completed.stderr)
        self.assertIn("SECRET-HIDDEN-SLIDE", raw)
        self.assertNotIn("VISIBLE-TITLE", raw)
        assert data is not None
        self.assertTrue(data["slides"][1]["included"])
        self.assertTrue(data["slides"][1]["hidden"])

        completed, data, raw, _ = self.run_inspector(None, "--include-hidden-slides")
        self.assertEqual(completed.returncode, 0, completed.stderr)
        self.assertIn("VISIBLE-TITLE", raw)
        self.assertIn("SECRET-HIDDEN-SLIDE", raw)
        assert data is not None
        self.assertEqual(data["presentation"]["included_slide_count"], 2)

    def test_group_transform_is_applied_before_off_slide_exclusion(self) -> None:
        base = self.root / "group-transform-base.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        group = slide.shapes.add_group_shape()
        visible = group.shapes.add_textbox(
            Inches(0.05), Inches(0.1), Inches(0.05), Inches(0.05)
        )
        visible.text = "VISIBLE-GROUP-CHILD"
        secret = group.shapes.add_textbox(
            Inches(0.9), Inches(0.1), Inches(0.05), Inches(0.05)
        )
        secret.text = "SECRET-TRANSFORMED-OFFSLIDE"
        presentation.save(base)

        with zipfile.ZipFile(base, "r") as archive:
            slide_xml = archive.read("ppt/slides/slide1.xml")
        root = ElementTree.fromstring(slide_xml)
        group_element = next(
            element for element in root.iter() if local_name(element.tag) == "grpSp"
        )
        group_properties = next(
            child for child in group_element if local_name(child.tag) == "grpSpPr"
        )
        transform = next(
            child for child in group_properties if local_name(child.tag) == "xfrm"
        )
        values = {
            "off": {"x": "-914400", "y": "0"},
            "ext": {"cx": "14630400", "cy": "914400"},
            "chOff": {"x": "0", "y": "0"},
            "chExt": {"cx": "914400", "cy": "914400"},
        }
        for child in transform:
            name = local_name(child.tag)
            if name in values:
                child.attrib.update(values[name])

        transformed = self.root / "group-transform.pptx"
        rewrite_zip(
            base,
            transformed,
            replacements={
                "ppt/slides/slide1.xml": ElementTree.tostring(
                    root, encoding="utf-8", xml_declaration=True
                )
            },
        )

        completed, _, raw, _ = self.run_inspector(transformed)
        self.assertEqual(completed.returncode, 0, completed.stderr)
        self.assertIn("VISIBLE-GROUP-CHILD", raw)
        self.assertNotIn("SECRET-TRANSFORMED-OFFSLIDE", raw)

        completed, _, raw, _ = self.run_inspector(
            transformed, "--include-off-slide-shapes"
        )
        self.assertEqual(completed.returncode, 0, completed.stderr)
        self.assertIn("VISIBLE-GROUP-CHILD", raw)
        self.assertIn("SECRET-TRANSFORMED-OFFSLIDE", raw)

    def test_rotated_polygon_outside_slide_is_excluded_even_when_bounds_overlap(self) -> None:
        rotated = self.root / "rotated-off-slide.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        textbox = slide.shapes.add_textbox(
            Inches(-1.848528),
            Inches(-1.848528),
            Inches(1.697056),
            Inches(1.697056),
        )
        textbox.text = "SECRET-ROTATED-FULLY-OFFSLIDE"
        textbox.rotation = 45
        presentation.save(rotated)

        completed, _, raw, _ = self.run_inspector(rotated)
        self.assertEqual(completed.returncode, 0, completed.stderr)
        self.assertNotIn("SECRET-ROTATED-FULLY-OFFSLIDE", raw)

        completed, _, raw, _ = self.run_inspector(
            rotated, "--include-off-slide-shapes"
        )
        self.assertEqual(completed.returncode, 0, completed.stderr)
        self.assertIn("SECRET-ROTATED-FULLY-OFFSLIDE", raw)

    def test_zero_area_connector_inside_slide_is_not_excluded(self) -> None:
        connector_file = self.root / "connector.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(1),
            Inches(1),
            Inches(3),
            Inches(1),
        )
        presentation.save(connector_file)

        completed, data, _, _ = self.run_inspector(connector_file)
        self.assertEqual(completed.returncode, 0, completed.stderr)
        assert data is not None
        connector_entry = next(
            shape
            for shape in data["slides"][0]["shapes"]
            if shape["shape_id"] == connector.shape_id
        )
        self.assertTrue(connector_entry["included"])
        self.assertFalse(connector_entry["fully_off_slide"])

    def test_json_is_deterministic_and_images_are_extracted_without_overwrite(self) -> None:
        image_directory = self.root / "assets"
        completed1, data1, raw1, _ = self.run_inspector(
            None, "--extract-images", str(image_directory)
        )
        completed2, data2, raw2, _ = self.run_inspector(
            None, "--extract-images", str(image_directory)
        )
        self.assertEqual(completed1.returncode, 0, completed1.stderr)
        self.assertEqual(completed2.returncode, 0, completed2.stderr)
        self.assertEqual(raw1, raw2)
        self.assertEqual(data1, data2)
        images = list(image_directory.iterdir())
        self.assertEqual(len(images), 1)
        self.assertGreater(images[0].stat().st_size, 0)

    def test_existing_image_with_different_bytes_is_never_overwritten(self) -> None:
        image_directory = self.root / "protected-assets"
        completed, data, _, _ = self.run_inspector(
            None, "--extract-images", str(image_directory)
        )
        self.assertEqual(completed.returncode, 0, completed.stderr)
        assert data is not None
        extracted_path = next(
            Path(shape["image"]["extracted_path"])
            for shape in iter_shapes(data["slides"][0]["shapes"])
            if shape.get("image") and shape["image"].get("extracted_path")
        )
        sentinel = b"EXISTING-IMAGE-MUST-NOT-BE-OVERWRITTEN"
        extracted_path.write_bytes(sentinel)

        completed, _, _, output = self.run_inspector(
            None, "--extract-images", str(image_directory)
        )
        self.assertEqual(completed.returncode, 2)
        self.assertIn("内容が異なります", completed.stderr)
        self.assertEqual(extracted_path.read_bytes(), sentinel)
        self.assertFalse(output.exists())

    def test_interrupted_image_write_leaves_no_partial_file(self) -> None:
        spec = importlib.util.spec_from_file_location("inspect_powerpoint", SCRIPT)
        assert spec is not None and spec.loader is not None
        module = importlib.util.module_from_spec(spec)
        sys.modules[spec.name] = module
        spec.loader.exec_module(module)
        image_stream = io.BytesIO()
        Image.new("RGB", (3, 2), (1, 2, 3)).save(image_stream, format="PNG")
        image_data = image_stream.getvalue()
        image_directory = self.root / "interrupted-assets"
        collector = module.ImageCollector(str(image_directory), 1024 * 1024)
        extracted_path, warning = collector.add(
            slide_number=1,
            shape_id=7,
            extension="png",
            content_type="image/png",
            data=image_data,
        )
        self.assertIsNotNone(extracted_path)
        self.assertIsNone(warning)

        with mock.patch.object(module.os, "fsync", side_effect=OSError("forced")):
            with self.assertRaises(OSError):
                collector.write_all()
        assert extracted_path is not None
        self.assertFalse(Path(extracted_path).exists())
        self.assertEqual(list(image_directory.iterdir()), [])

    def test_json_output_requires_json_and_force_for_overwrite(self) -> None:
        non_json = self.root / "inspection.txt"
        completed = subprocess.run(
            [
                sys.executable,
                str(SCRIPT),
                str(self.presentation),
                "--output",
                str(non_json),
            ],
            capture_output=True,
            text=True,
            check=False,
        )
        self.assertEqual(completed.returncode, 2)
        self.assertIn("拡張子.json", completed.stderr)
        self.assertFalse(non_json.exists())

        fresh = self.root / "fresh.JSON"
        completed = subprocess.run(
            [
                sys.executable,
                str(SCRIPT),
                str(self.presentation),
                "--output",
                str(fresh),
            ],
            capture_output=True,
            text=True,
            check=False,
        )
        self.assertEqual(completed.returncode, 0, completed.stderr)
        self.assertEqual(json.loads(fresh.read_text(encoding="utf-8"))["schema_version"], "1.1")
        self.assertEqual(list(self.root.glob(f".{fresh.name}.*.tmp")), [])

        existing = self.root / "existing.json"
        sentinel = '{"sentinel": true}\n'
        existing.write_text(sentinel, encoding="utf-8")
        command = [
            sys.executable,
            str(SCRIPT),
            str(self.presentation),
            "--output",
            str(existing),
        ]
        completed = subprocess.run(command, capture_output=True, text=True, check=False)
        self.assertEqual(completed.returncode, 2)
        self.assertIn("--force", completed.stderr)
        self.assertEqual(existing.read_text(encoding="utf-8"), sentinel)

        completed = subprocess.run(
            [*command, "--force"], capture_output=True, text=True, check=False
        )
        self.assertEqual(completed.returncode, 0, completed.stderr)
        self.assertEqual(json.loads(existing.read_text(encoding="utf-8"))["schema_version"], "1.1")
        self.assertEqual(list(self.root.glob(f".{existing.name}.*.tmp")), [])

    def test_force_replaces_output_symlink_but_never_the_input_target(self) -> None:
        protected = self.root / "protected-target.txt"
        protected.write_text("DO-NOT-CHANGE", encoding="utf-8")
        output_link = self.root / "linked-output.json"
        try:
            output_link.symlink_to(protected)
        except (OSError, NotImplementedError):
            self.skipTest("この環境ではシンボリックリンクを作成できません。")

        command = [
            sys.executable,
            str(SCRIPT),
            str(self.presentation),
            "--output",
            str(output_link),
        ]
        completed = subprocess.run(command, capture_output=True, text=True, check=False)
        self.assertEqual(completed.returncode, 2)
        self.assertTrue(output_link.is_symlink())
        self.assertEqual(protected.read_text(encoding="utf-8"), "DO-NOT-CHANGE")

        completed = subprocess.run(
            [*command, "--force"], capture_output=True, text=True, check=False
        )
        self.assertEqual(completed.returncode, 0, completed.stderr)
        self.assertFalse(output_link.is_symlink())
        self.assertEqual(protected.read_text(encoding="utf-8"), "DO-NOT-CHANGE")
        self.assertEqual(json.loads(output_link.read_text(encoding="utf-8"))["schema_version"], "1.1")

        input_link = self.root / "input-alias.json"
        input_link.symlink_to(self.presentation)
        completed = subprocess.run(
            [
                sys.executable,
                str(SCRIPT),
                str(self.presentation),
                "--output",
                str(input_link),
                "--force",
            ],
            capture_output=True,
            text=True,
            check=False,
        )
        self.assertEqual(completed.returncode, 2)
        self.assertIn("同じパス", completed.stderr)
        self.assertTrue(input_link.is_symlink())

        input_hardlink = self.root / "input-hardlink.json"
        try:
            input_hardlink.hardlink_to(self.presentation)
        except OSError:
            self.skipTest("この環境ではハードリンクを作成できません。")
        before = hashlib.sha256(self.presentation.read_bytes()).hexdigest()
        completed = subprocess.run(
            [
                sys.executable,
                str(SCRIPT),
                str(self.presentation),
                "--output",
                str(input_hardlink),
                "--force",
            ],
            capture_output=True,
            text=True,
            check=False,
        )
        self.assertEqual(completed.returncode, 2)
        self.assertIn("同じファイル", completed.stderr)
        self.assertEqual(
            hashlib.sha256(self.presentation.read_bytes()).hexdigest(), before
        )

    def test_json_output_cannot_collide_with_an_extracted_image(self) -> None:
        spec = importlib.util.spec_from_file_location("inspect_powerpoint", SCRIPT)
        assert spec is not None and spec.loader is not None
        module = importlib.util.module_from_spec(spec)
        sys.modules[spec.name] = module
        spec.loader.exec_module(module)
        collision_path = self.root / "collision-assets" / "image.json"
        payload = b"image"
        collector = module.ImageCollector(str(collision_path.parent), 1024)
        collector.pending.append(
            module.PendingImage(
                collision_path,
                payload,
                hashlib.sha256(payload).hexdigest(),
            )
        )
        with self.assertRaisesRegex(module.InspectionError, "抽出画像の出力先と衝突"):
            collector.ensure_output_does_not_collide(str(collision_path))
        self.assertFalse(collision_path.exists())

    def test_images_created_in_a_failed_run_are_removed(self) -> None:
        spec = importlib.util.spec_from_file_location("inspect_powerpoint", SCRIPT)
        assert spec is not None and spec.loader is not None
        module = importlib.util.module_from_spec(spec)
        sys.modules[spec.name] = module
        spec.loader.exec_module(module)
        output_directory = self.root / "failed-output-assets"
        output = self.root / "forced-write-failure.json"
        with mock.patch.object(module, "write_json", side_effect=OSError("forced")):
            with contextlib.redirect_stderr(io.StringIO()):
                return_code = module.main(
                    [
                        str(self.presentation),
                        "--output",
                        str(output),
                        "--extract-images",
                        str(output_directory),
                    ]
                )
        self.assertEqual(return_code, 2)
        self.assertTrue(output_directory.is_dir())
        self.assertEqual(list(output_directory.iterdir()), [])
        self.assertFalse(output.exists())

    def test_limits_stop_instead_of_silently_truncating(self) -> None:
        for option, value in (
            ("--max-slides", "1"),
            ("--max-shapes", "1"),
            ("--max-text-chars", "5"),
            ("--max-table-cells", "1"),
            ("--max-chart-points", "1"),
        ):
            completed, data, _, output = self.run_inspector(None, option, value)
            self.assertEqual(completed.returncode, 2, (option, completed.stderr))
            self.assertIsNone(data)
            self.assertFalse(output.exists())
            self.assertIn("上限", completed.stderr)

    def test_unsupported_and_malformed_files_are_rejected(self) -> None:
        legacy = self.root / "legacy.ppt"
        legacy.write_bytes(b"not a presentation")
        completed, data, _, _ = self.run_inspector(legacy)
        self.assertEqual(completed.returncode, 2)
        self.assertIsNone(data)
        self.assertIn(".pptxへ変換", completed.stderr)

        malformed = self.root / "malformed.pptx"
        malformed.write_bytes(b"not a zip")
        completed, data, _, _ = self.run_inspector(malformed)
        self.assertEqual(completed.returncode, 2)
        self.assertIsNone(data)
        self.assertIn("ZIP形式", completed.stderr)

    def test_utf16_dtd_with_disguised_extension_is_rejected(self) -> None:
        malicious = self.root / "dtd.pptx"
        payload = "<?xml version='1.0'?><!DOCTYPE x [<!ENTITY e 'x'>]><x>&e;</x>".encode(
            "utf-16"
        )
        rewrite_zip(
            self.presentation,
            malicious,
            additions=[("ppt/media/disguised.BIN", payload)],
        )
        completed, data, _, _ = self.run_inspector(malicious)
        self.assertEqual(completed.returncode, 2)
        self.assertIsNone(data)
        self.assertIn("DTDまたはXMLエンティティ", completed.stderr)

    def test_unsafe_archive_names_case_collisions_and_symlinks_are_rejected(self) -> None:
        traversal = self.root / "traversal.pptx"
        rewrite_zip(self.presentation, traversal, additions=[("../outside.xml", b"<x/>")])
        completed, _, _, _ = self.run_inspector(traversal)
        self.assertEqual(completed.returncode, 2)
        self.assertIn("ZIP部品名", completed.stderr)

        encoded_traversal = self.root / "encoded-traversal.pptx"
        rewrite_zip(
            self.presentation,
            encoded_traversal,
            additions=[("ppt/media/%252e%252e/secret.bin", b"x")],
        )
        completed, _, _, _ = self.run_inspector(encoded_traversal)
        self.assertEqual(completed.returncode, 2)
        self.assertIn("エンコード後パス", completed.stderr)

        collision = self.root / "collision.pptx"
        rewrite_zip(
            self.presentation,
            collision,
            additions=[("PPT/PRESENTATION.XML", b"<x/>")],
        )
        completed, _, _, _ = self.run_inspector(collision)
        self.assertEqual(completed.returncode, 2)
        self.assertIn("重複または大小文字衝突", completed.stderr)

        symlink = self.root / "symlink.pptx"
        link_info = zipfile.ZipInfo("ppt/media/link.png")
        link_info.create_system = 3
        link_info.external_attr = (stat.S_IFLNK | 0o777) << 16
        rewrite_zip(self.presentation, symlink, additions=[(link_info, b"../../outside")])
        completed, _, _, _ = self.run_inspector(symlink)
        self.assertEqual(completed.returncode, 2)
        self.assertIn("シンボリックリンク", completed.stderr)

    def test_relationship_escaping_package_is_rejected(self) -> None:
        with zipfile.ZipFile(self.presentation, "r") as archive:
            root_relationships = archive.read("_rels/.rels")
        root = ElementTree.fromstring(root_relationships)
        ElementTree.SubElement(
            root,
            f"{{{REL_NS}}}Relationship",
            {
                "Id": "rIdTraversal",
                "Type": "http://example.invalid/relationship/test",
                "Target": "../../outside.xml",
            },
        )
        malicious = self.root / "relationship.pptx"
        rewrite_zip(
            self.presentation,
            malicious,
            replacements={
                "_rels/.rels": ElementTree.tostring(root, encoding="utf-8", xml_declaration=True)
            },
        )
        completed, data, _, _ = self.run_inspector(malicious)
        self.assertEqual(completed.returncode, 2)
        self.assertIsNone(data)
        self.assertIn("パッケージ外", completed.stderr)

        encoded_root = ElementTree.fromstring(root_relationships)
        ElementTree.SubElement(
            encoded_root,
            f"{{{REL_NS}}}Relationship",
            {
                "Id": "rIdEncodedTraversal",
                "Type": "http://example.invalid/relationship/test",
                "Target": "%252e%252e/%252e%252e/outside.xml",
            },
        )
        encoded = self.root / "encoded-relationship.pptx"
        rewrite_zip(
            self.presentation,
            encoded,
            replacements={
                "_rels/.rels": ElementTree.tostring(
                    encoded_root, encoding="utf-8", xml_declaration=True
                )
            },
        )
        completed, data, _, _ = self.run_inspector(encoded)
        self.assertEqual(completed.returncode, 2)
        self.assertIsNone(data)
        self.assertIn("パッケージ外", completed.stderr)

    def test_vector_images_are_not_written_without_sanitization(self) -> None:
        spec = importlib.util.spec_from_file_location("inspect_powerpoint", SCRIPT)
        assert spec is not None and spec.loader is not None
        module = importlib.util.module_from_spec(spec)
        sys.modules[spec.name] = module
        spec.loader.exec_module(module)
        collector = module.ImageCollector(str(self.root / "vector-assets"), 1024)
        extracted, warning = collector.add(
            slide_number=1,
            shape_id=99,
            extension="svg",
            data=b"<svg xmlns='http://www.w3.org/2000/svg'><script/></svg>",
        )
        self.assertIsNone(extracted)
        self.assertIn("自動抽出していません", warning)
        self.assertEqual(collector.pending, [])
        collector.write_all()
        self.assertFalse((self.root / "vector-assets").exists())

    def test_svg_picture_keeps_ooxml_metadata_when_python_pptx_cannot_decode_it(self) -> None:
        svg_data = (
            b"<svg xmlns='http://www.w3.org/2000/svg'>"
            b"<script>DO-NOT-EXECUTE</script><rect width='10' height='10'/></svg>"
        )
        vector = self.root / "vector-picture.pptx"
        retarget_first_picture(
            self.presentation,
            vector,
            filename="unsafe.svg",
            content_type="image/svg+xml",
            data=svg_data,
        )
        image_directory = self.root / "vector-picture-assets"

        completed, data, raw, _ = self.run_inspector(
            vector, "--extract-images", str(image_directory)
        )
        self.assertEqual(completed.returncode, 0, completed.stderr)
        assert data is not None
        image = next(
            shape["image"]
            for shape in iter_shapes(data["slides"][0]["shapes"])
            if shape.get("image")
        )
        self.assertEqual(image["filename"], "unsafe.svg")
        self.assertEqual(image["content_type"], "image/svg+xml")
        self.assertEqual(image["extension"], "svg")
        self.assertEqual(image["bytes"], len(svg_data))
        self.assertEqual(image["sha256"], hashlib.sha256(svg_data).hexdigest())
        self.assertFalse(image["extractable"])
        self.assertIsNone(image["extracted_path"])
        self.assertIn("自動抽出していません", image["not_extracted_reason"])
        self.assertIn("OOXML部品", image["inspection_warning"])
        self.assertFalse(image_directory.exists())
        self.assertNotIn("DO-NOT-EXECUTE", raw)

    def test_svg_extension_part_is_preferred_over_png_fallback(self) -> None:
        relationship_part = "ppt/slides/_rels/slide1.xml.rels"
        slide_part = "ppt/slides/slide1.xml"
        with zipfile.ZipFile(self.presentation, "r") as archive:
            relationships_root = ElementTree.fromstring(
                archive.read(relationship_part)
            )
            slide_root = ElementTree.fromstring(archive.read(slide_part))
            content_types_root = ElementTree.fromstring(
                archive.read("[Content_Types].xml")
            )

        svg_relationship_id = "rIdSvgOriginal"
        ElementTree.SubElement(
            relationships_root,
            f"{{{REL_NS}}}Relationship",
            {
                "Id": svg_relationship_id,
                "Type": IMAGE_REL_TYPE,
                "Target": "../media/original.svg",
            },
        )
        blip = next(element for element in slide_root.iter() if local_name(element.tag) == "blip")
        extension_list = ElementTree.SubElement(blip, f"{{{DRAWING_NS}}}extLst")
        extension = ElementTree.SubElement(
            extension_list,
            f"{{{DRAWING_NS}}}ext",
            {"uri": "{96DAC541-7B7A-43D3-8B79-37D633B846F1}"},
        )
        ElementTree.SubElement(
            extension,
            f"{{{SVG_NS}}}svgBlip",
            {f"{{{OFFICE_REL_NS}}}embed": svg_relationship_id},
        )
        ElementTree.SubElement(
            content_types_root,
            f"{{{CONTENT_TYPES_NS}}}Default",
            {"Extension": "svg", "ContentType": "image/svg+xml"},
        )
        svg_data = b"<svg xmlns='http://www.w3.org/2000/svg'><circle r='4'/></svg>"
        dual = self.root / "svg-with-png-fallback.pptx"
        rewrite_zip(
            self.presentation,
            dual,
            replacements={
                relationship_part: ElementTree.tostring(
                    relationships_root, encoding="utf-8", xml_declaration=True
                ),
                slide_part: ElementTree.tostring(
                    slide_root, encoding="utf-8", xml_declaration=True
                ),
                "[Content_Types].xml": ElementTree.tostring(
                    content_types_root, encoding="utf-8", xml_declaration=True
                ),
            },
            additions=[("ppt/media/original.svg", svg_data)],
        )
        image_directory = self.root / "dual-svg-assets"

        completed, data, _, _ = self.run_inspector(
            dual, "--extract-images", str(image_directory)
        )
        self.assertEqual(completed.returncode, 0, completed.stderr)
        assert data is not None
        image = next(
            shape["image"]
            for shape in iter_shapes(data["slides"][0]["shapes"])
            if shape.get("image")
        )
        self.assertEqual(image["filename"], "original.svg")
        self.assertEqual(image["content_type"], "image/svg+xml")
        self.assertEqual(image["sha256"], hashlib.sha256(svg_data).hexdigest())
        self.assertFalse(image["extractable"])
        self.assertIsNone(image["extracted_path"])
        self.assertIsNone(image["pixel_size"])
        self.assertIn("PNGフォールバックより優先", image["resolution_warning"])
        self.assertIn("フォールバック画像", image["inspection_warning"])
        self.assertFalse(image_directory.exists())

    def test_webp_picture_is_resolved_from_ooxml_and_extracted_when_requested(self) -> None:
        webp_stream = io.BytesIO()
        Image.new("RGB", (7, 5), (20, 120, 220)).save(webp_stream, format="WEBP")
        webp_data = webp_stream.getvalue()
        webp = self.root / "webp-picture.pptx"
        retarget_first_picture(
            self.presentation,
            webp,
            filename="fallback.webp",
            content_type="image/webp",
            data=webp_data,
        )
        image_directory = self.root / "webp-picture-assets"

        completed, data, _, _ = self.run_inspector(
            webp, "--extract-images", str(image_directory)
        )
        self.assertEqual(completed.returncode, 0, completed.stderr)
        assert data is not None
        image = next(
            shape["image"]
            for shape in iter_shapes(data["slides"][0]["shapes"])
            if shape.get("image")
        )
        self.assertEqual(image["filename"], "fallback.webp")
        self.assertEqual(image["content_type"], "image/webp")
        self.assertEqual(image["extension"], "webp")
        self.assertEqual(image["bytes"], len(webp_data))
        self.assertEqual(image["sha256"], hashlib.sha256(webp_data).hexdigest())
        self.assertTrue(image["extractable"])
        self.assertIsNone(image["not_extracted_reason"])
        extracted_path = Path(image["extracted_path"])
        self.assertEqual(extracted_path.read_bytes(), webp_data)
        self.assertEqual(extracted_path.suffix, ".webp")

    def test_external_picture_relationship_is_reported_without_opening_target(self) -> None:
        external = self.root / "external-picture.pptx"
        secret_target = "https://example.invalid/SECRET-EXTERNAL-IMAGE.svg"
        retarget_first_picture(
            self.presentation,
            external,
            external_target=secret_target,
        )
        image_directory = self.root / "external-picture-assets"

        completed, data, raw, _ = self.run_inspector(
            external, "--extract-images", str(image_directory)
        )
        self.assertEqual(completed.returncode, 0, completed.stderr)
        assert data is not None
        image = next(
            shape["image"]
            for shape in iter_shapes(data["slides"][0]["shapes"])
            if shape.get("image")
        )
        for field in ("filename", "content_type", "extension", "bytes", "sha256"):
            self.assertIsNone(image[field])
        self.assertFalse(image["extractable"])
        self.assertIsNone(image["extracted_path"])
        self.assertIn("外部画像Relationship", image["not_extracted_reason"])
        self.assertNotIn("SECRET-EXTERNAL-IMAGE", raw)
        self.assertFalse(image_directory.exists())

    def test_picture_relationship_with_unsafe_internal_target_is_rejected(self) -> None:
        relationship_part = "ppt/slides/_rels/slide1.xml.rels"
        with zipfile.ZipFile(self.presentation, "r") as archive:
            relationships_root = ElementTree.fromstring(
                archive.read(relationship_part)
            )
        image_relationship = next(
            element
            for element in relationships_root
            if element.attrib.get("Type") == IMAGE_REL_TYPE
        )
        image_relationship.set("Target", "../../../../outside.png")
        unsafe = self.root / "unsafe-picture-target.pptx"
        rewrite_zip(
            self.presentation,
            unsafe,
            replacements={
                relationship_part: ElementTree.tostring(
                    relationships_root, encoding="utf-8", xml_declaration=True
                )
            },
        )

        completed, data, _, output = self.run_inspector(unsafe)
        self.assertEqual(completed.returncode, 2)
        self.assertIn("パッケージ外", completed.stderr)
        self.assertIsNone(data)
        self.assertFalse(output.exists())

    def test_macro_enabled_package_is_read_only_and_vba_is_inventory_only(self) -> None:
        with zipfile.ZipFile(self.presentation, "r") as archive:
            content_types = archive.read("[Content_Types].xml")
            presentation_relationships = archive.read("ppt/_rels/presentation.xml.rels")

        content_root = ElementTree.fromstring(content_types)
        for element in content_root:
            if (
                local_name(element.tag) == "Override"
                and element.attrib.get("PartName") == "/ppt/presentation.xml"
            ):
                element.set(
                    "ContentType",
                    "application/vnd.ms-powerpoint.presentation.macroEnabled.main+xml",
                )
        ElementTree.SubElement(
            content_root,
            f"{{{CONTENT_TYPES_NS}}}Override",
            {
                "PartName": "/ppt/vbaProject.bin",
                "ContentType": "application/vnd.ms-office.vbaProject",
            },
        )

        relationships_root = ElementTree.fromstring(presentation_relationships)
        ElementTree.SubElement(
            relationships_root,
            f"{{{REL_NS}}}Relationship",
            {
                "Id": "rIdVba",
                "Type": "http://schemas.microsoft.com/office/2006/relationships/vbaProject",
                "Target": "vbaProject.bin",
            },
        )

        macro = self.root / "sample.pptm"
        rewrite_zip(
            self.presentation,
            macro,
            replacements={
                "[Content_Types].xml": ElementTree.tostring(
                    content_root, encoding="utf-8", xml_declaration=True
                ),
                "ppt/_rels/presentation.xml.rels": ElementTree.tostring(
                    relationships_root, encoding="utf-8", xml_declaration=True
                ),
            },
            additions=[("ppt/vbaProject.bin", b"NOT-EXECUTABLE-TEST-VBA")],
        )
        before = hashlib.sha256(macro.read_bytes()).hexdigest()
        completed, data, raw, _ = self.run_inspector(macro)
        after = hashlib.sha256(macro.read_bytes()).hexdigest()
        self.assertEqual(completed.returncode, 0, completed.stderr)
        self.assertEqual(before, after)
        assert data is not None
        self.assertTrue(data["presentation"]["archive"]["has_vba"])
        self.assertNotIn("NOT-EXECUTABLE-TEST-VBA", raw)

    def test_vba_is_detected_by_relationship_and_content_type_with_renamed_part(self) -> None:
        with zipfile.ZipFile(self.presentation, "r") as archive:
            content_types = archive.read("[Content_Types].xml")
            presentation_relationships = archive.read("ppt/_rels/presentation.xml.rels")

        content_root = ElementTree.fromstring(content_types)
        ElementTree.SubElement(
            content_root,
            f"{{{CONTENT_TYPES_NS}}}Override",
            {
                "PartName": "/ppt/customPayload.bin",
                "ContentType": "application/vnd.ms-office.vbaProject",
            },
        )
        relationships_root = ElementTree.fromstring(presentation_relationships)
        ElementTree.SubElement(
            relationships_root,
            f"{{{REL_NS}}}Relationship",
            {
                "Id": "rIdRenamedVba",
                "Type": "http://schemas.microsoft.com/office/2006/relationships/vbaProject",
                "Target": "customPayload.bin",
            },
        )
        renamed = self.root / "renamed-vba.pptx"
        rewrite_zip(
            self.presentation,
            renamed,
            replacements={
                "[Content_Types].xml": ElementTree.tostring(
                    content_root, encoding="utf-8", xml_declaration=True
                ),
                "ppt/_rels/presentation.xml.rels": ElementTree.tostring(
                    relationships_root, encoding="utf-8", xml_declaration=True
                ),
            },
            additions=[("ppt/customPayload.bin", b"RENAMED-TEST-VBA")],
        )
        completed, data, raw, _ = self.run_inspector(renamed)
        self.assertEqual(completed.returncode, 0, completed.stderr)
        assert data is not None
        self.assertTrue(data["presentation"]["archive"]["has_vba"])
        self.assertNotIn("RENAMED-TEST-VBA", raw)

    def test_invalid_slide_and_input_output_collision_are_rejected(self) -> None:
        completed, data, _, _ = self.run_inspector(None, "--slide", "99")
        self.assertEqual(completed.returncode, 2)
        self.assertIsNone(data)
        self.assertIn("存在しないスライド番号", completed.stderr)

        completed = subprocess.run(
            [
                sys.executable,
                str(SCRIPT),
                str(self.presentation),
                "--output",
                str(self.presentation),
            ],
            capture_output=True,
            text=True,
            check=False,
        )
        self.assertEqual(completed.returncode, 2)
        self.assertIn("同じパス", completed.stderr)


if __name__ == "__main__":
    unittest.main()
