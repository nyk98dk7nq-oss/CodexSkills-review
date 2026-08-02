#!/usr/bin/env python3
"""PowerPointを安全に読み取り、Markdown変換用の決定的なJSONを生成する。"""

from __future__ import annotations

import argparse
import hashlib
import io
import json
import math
import os
import posixpath
import re
import stat
import sys
import tempfile
import unicodedata
import zipfile
import xml.etree.ElementTree as ElementTree
from dataclasses import dataclass, field
from datetime import date, datetime, time
from pathlib import Path, PurePosixPath
from typing import Any, Iterable
from urllib.parse import unquote, urlsplit

if sys.version_info < (3, 10):  # pragma: no cover - 実行環境依存
    print("エラー: inspect_powerpoint.pyにはPython 3.10以上が必要です。", file=sys.stderr)
    raise SystemExit(3)

try:
    import pptx
    from pptx import Presentation
except ImportError as exc:  # pragma: no cover - 実行環境依存
    print(
        "エラー: inspect_powerpoint.pyにはpython-pptx 1.0.2以上2.0未満が必要です。",
        file=sys.stderr,
    )
    raise SystemExit(3) from exc

_PPTX_VERSION = re.match(r"^(\d+)\.(\d+)\.(\d+)", pptx.__version__)
if _PPTX_VERSION is None:
    print(f"エラー: python-pptxのバージョンを判定できません: {pptx.__version__}", file=sys.stderr)
    raise SystemExit(3)
_PPTX_VERSION_TUPLE = tuple(map(int, _PPTX_VERSION.groups()))
if _PPTX_VERSION_TUPLE < (1, 0, 2) or _PPTX_VERSION_TUPLE >= (2, 0, 0):
    print(
        "エラー: python-pptx 1.0.2以上2.0未満が必要です。"
        f"現在: {pptx.__version__}",
        file=sys.stderr,
    )
    raise SystemExit(3)


SCHEMA_VERSION = "1.1"
SUPPORTED_SUFFIXES = {".pptx", ".pptm"}
LEGACY_SUFFIXES = {".ppt", ".pps", ".pot"}
OTHER_OPENXML_SUFFIXES = {".ppsx", ".ppsm", ".potx", ".potm"}

DEFAULT_MAX_FILE_MB = 50
DEFAULT_MAX_UNCOMPRESSED_MB = 500
DEFAULT_MAX_SLIDES = 500
DEFAULT_MAX_SHAPES = 10_000
DEFAULT_MAX_TABLE_CELLS = 50_000
DEFAULT_MAX_CHART_POINTS = 100_000
DEFAULT_MAX_TEXT_CHARS = 2_000_000
DEFAULT_MAX_EXTRACTED_IMAGES_MB = 200

MAX_ARCHIVE_PARTS = 20_000
MAX_SUSPICIOUS_COMPRESSION_RATIO = 1_000
MAX_XML_PART_BYTES = 20 * 1024 * 1024
MAX_RELATIONSHIPS_BYTES = 10 * 1024 * 1024
EMU_PER_INCH = 914_400

IMAGE_SUFFIXES = {
    ".apng",
    ".bmp",
    ".emf",
    ".gif",
    ".heic",
    ".heif",
    ".jpeg",
    ".jpg",
    ".png",
    ".svg",
    ".tif",
    ".tiff",
    ".webp",
    ".wmf",
}
SAFE_EXTRACTABLE_IMAGE_SUFFIXES = {
    "apng",
    "bmp",
    "gif",
    "jpeg",
    "jpg",
    "png",
    "tif",
    "tiff",
    "webp",
}
SAFE_EXTRACTABLE_IMAGE_CONTENT_TYPES = {
    "apng": {"image/apng", "image/png"},
    "bmp": {"image/bmp", "image/x-ms-bmp"},
    "gif": {"image/gif"},
    "jpeg": {"image/jpeg"},
    "jpg": {"image/jpeg"},
    "png": {"image/png"},
    "tif": {"image/tiff"},
    "tiff": {"image/tiff"},
    "webp": {"image/webp"},
}
IMAGE_RELATIONSHIP_TYPES = {
    "http://purl.oclc.org/ooxml/officeDocument/relationships/image",
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image",
}
AUDIO_VIDEO_SUFFIXES = {
    ".aac",
    ".avi",
    ".m4a",
    ".m4v",
    ".mid",
    ".midi",
    ".mov",
    ".mp3",
    ".mp4",
    ".mpeg",
    ".mpg",
    ".wav",
    ".wma",
    ".wmv",
}
BIDI_AND_DIRECTIONAL_CONTROLS = {
    0x061C,
    0x200E,
    0x200F,
    0x202A,
    0x202B,
    0x202C,
    0x202D,
    0x202E,
    0x2066,
    0x2067,
    0x2068,
    0x2069,
}

# Relationship Typeは非表示スライドなど除外対象の部品にも属し得る。任意の末尾
# 文字列をJSONのキーへ出すと除外内容が漏れるため、既知の固定ラベルだけを公開し、
# それ以外はotherへ集約する。
SAFE_RELATIONSHIP_TYPE_LABELS = {
    "activeXControl",
    "activeXControlBinary",
    "audio",
    "chart",
    "chartUserShapes",
    "commentAuthors",
    "comments",
    "customProperties",
    "customXml",
    "diagramColors",
    "diagramData",
    "diagramLayout",
    "diagramQuickStyle",
    "externalLink",
    "font",
    "handoutMaster",
    "hyperlink",
    "image",
    "media",
    "notesMaster",
    "notesSlide",
    "officeDocument",
    "oleObject",
    "package",
    "presProps",
    "slide",
    "slideLayout",
    "slideMaster",
    "tableStyles",
    "tags",
    "theme",
    "userShapes",
    "vbaProject",
    "video",
    "viewProps",
}


class InspectionError(RuntimeError):
    """ユーザーが対処できる検査エラー。"""


@dataclass
class Limits:
    max_slides: int
    max_shapes: int
    max_table_cells: int
    max_chart_points: int
    max_text_chars: int
    shape_count: int = 0
    table_cell_count: int = 0
    chart_point_count: int = 0
    text_char_count: int = 0

    def add_shapes(self, count: int) -> None:
        self.shape_count += count
        if self.shape_count > self.max_shapes:
            raise InspectionError(
                f"抽出対象の図形数が上限{self.max_shapes}を超えています。"
                "--slideで対象を絞るか、安全性を確認してから上限を変更してください。"
            )

    def add_table_cells(self, count: int) -> None:
        self.table_cell_count += count
        if self.table_cell_count > self.max_table_cells:
            raise InspectionError(
                f"抽出対象の表セル数が上限{self.max_table_cells}を超えています。"
                "--slideで対象を絞るか、安全性を確認してから上限を変更してください。"
            )

    def add_chart_points(self, count: int) -> None:
        self.chart_point_count += count
        if self.chart_point_count > self.max_chart_points:
            raise InspectionError(
                f"抽出対象のグラフ点数が上限{self.max_chart_points}を超えています。"
                "--slideで対象を絞るか、安全性を確認してから上限を変更してください。"
            )

    def add_text(self, value: str) -> None:
        self.text_char_count += len(value)
        if self.text_char_count > self.max_text_chars:
            raise InspectionError(
                f"抽出対象の文字数が上限{self.max_text_chars}を超えています。"
                "--slideで対象を絞るか、安全性を確認してから上限を変更してください。"
            )

    def as_dict(self) -> dict[str, int]:
        return {
            "max_slides": self.max_slides,
            "max_shapes": self.max_shapes,
            "max_table_cells": self.max_table_cells,
            "max_chart_points": self.max_chart_points,
            "max_text_chars": self.max_text_chars,
        }


@dataclass(frozen=True)
class AffineTransform:
    """親座標からスライド座標への2次元アフィン変換。"""

    a: float = 1.0
    b: float = 0.0
    c: float = 0.0
    d: float = 1.0
    e: float = 0.0
    f: float = 0.0

    def point(self, x: float, y: float) -> tuple[float, float]:
        return (
            self.a * x + self.c * y + self.e,
            self.b * x + self.d * y + self.f,
        )

    def compose(self, inner: "AffineTransform") -> "AffineTransform":
        """self(inner(point))を返す。"""

        return AffineTransform(
            a=self.a * inner.a + self.c * inner.b,
            b=self.b * inner.a + self.d * inner.b,
            c=self.a * inner.c + self.c * inner.d,
            d=self.b * inner.c + self.d * inner.d,
            e=self.a * inner.e + self.c * inner.f + self.e,
            f=self.b * inner.e + self.d * inner.f + self.f,
        )


@dataclass
class PendingImage:
    path: Path
    data: bytes
    sha256: str


@dataclass(frozen=True)
class PackageImage:
    filename: str | None
    content_type: str | None
    extension: str | None
    data: bytes | None
    resolution_warning: str | None = None
    decode_with_python_pptx: bool = True


def raster_signature_matches(extension: str, data: bytes) -> bool:
    if extension in {"png", "apng"}:
        return data.startswith(b"\x89PNG\r\n\x1a\n")
    if extension in {"jpeg", "jpg"}:
        return data.startswith(b"\xff\xd8\xff")
    if extension == "gif":
        return data.startswith((b"GIF87a", b"GIF89a"))
    if extension == "bmp":
        return data.startswith(b"BM")
    if extension in {"tif", "tiff"}:
        return data.startswith((b"II*\x00", b"MM\x00*"))
    if extension == "webp":
        return len(data) >= 12 and data.startswith(b"RIFF") and data[8:12] == b"WEBP"
    return False


def image_extraction_decision(
    extension: str | None,
    content_type: str | None,
    data: bytes | None,
) -> tuple[bool, str | None]:
    suffix = re.sub(r"[^a-z0-9]", "", (extension or "").casefold())
    if suffix not in SAFE_EXTRACTABLE_IMAGE_SUFFIXES:
        return (
            False,
            "SVG、EMF、WMFなどのベクター画像または未認識形式は、"
            "安全なサニタイズを行えないため自動抽出していません。",
        )
    normalized_content_type = (content_type or "").split(";", 1)[0].strip().casefold()
    if normalized_content_type not in SAFE_EXTRACTABLE_IMAGE_CONTENT_TYPES[suffix]:
        return False, "拡張子とContent-Typeが安全なラスター画像として一致しません。"
    if data is None or not raster_signature_matches(suffix, data):
        return False, "画像のバイト列が拡張子で示された安全なラスター形式と一致しません。"
    return True, None


@dataclass
class ImageCollector:
    requested_directory: str | None
    max_bytes: int
    pending: list[PendingImage] = field(default_factory=list)
    total_bytes: int = 0

    @property
    def enabled(self) -> bool:
        return self.requested_directory is not None

    def add(
        self,
        *,
        slide_number: int,
        shape_id: int,
        extension: str,
        content_type: str | None = None,
        data: bytes,
    ) -> tuple[str | None, str | None]:
        suffix = re.sub(r"[^a-z0-9]", "", extension.casefold()) or "bin"
        extractable, rejection_reason = image_extraction_decision(
            suffix, content_type, data
        )
        if not extractable:
            return None, rejection_reason
        if not self.enabled:
            return None, "--extract-imagesが指定されていないため抽出していません。"
        self.total_bytes += len(data)
        if self.total_bytes > self.max_bytes:
            raise InspectionError(
                "抽出画像の合計サイズが上限を超えています。"
                "対象スライドを絞るか、安全性を確認してから上限を変更してください。"
            )
        assert self.requested_directory is not None
        directory = Path(self.requested_directory)
        digest = hashlib.sha256(data).hexdigest()
        filename = f"slide-{slide_number:04d}-shape-{shape_id}-{digest[:12]}.{suffix}"
        target = directory / filename
        self.pending.append(PendingImage(target, data, digest))
        return target.as_posix(), None

    def ensure_output_does_not_collide(self, output: str | None) -> None:
        if output is None:
            return
        output_path = Path(output).resolve()
        if any(item.path.resolve() == output_path for item in self.pending):
            raise InspectionError(
                "JSONの出力先が抽出画像の出力先と衝突しています。"
                "--outputまたは--extract-imagesを変更してください。"
            )

    def write_all(self) -> list[Path]:
        if not self.pending:
            return []
        assert self.requested_directory is not None
        directory = Path(self.requested_directory)
        if directory.exists() and directory.is_symlink():
            raise InspectionError("画像の出力先にシンボリックリンクは指定できません。")
        if directory.exists() and not directory.is_dir():
            raise InspectionError("画像の出力先はディレクトリである必要があります。")
        directory.mkdir(parents=True, exist_ok=True)

        def verify_existing(item: PendingImage) -> None:
            if item.path.is_symlink() or not item.path.is_file():
                raise InspectionError(
                    f"画像の出力先が通常ファイルではありません: {item.path}"
                )
            current = hashlib.sha256(item.path.read_bytes()).hexdigest()
            if current != item.sha256:
                raise InspectionError(
                    f"同名の画像ファイルが既に存在し、内容が異なります: {item.path}"
                )

        created: list[Path] = []
        try:
            for item in self.pending:
                if os.path.lexists(item.path):
                    verify_existing(item)
                    continue
                descriptor, temporary_name = tempfile.mkstemp(
                    prefix=f".{item.path.name}.", suffix=".tmp", dir=directory
                )
                temporary_path = Path(temporary_name)
                try:
                    with os.fdopen(descriptor, "wb") as stream:
                        stream.write(item.data)
                        stream.flush()
                        os.fsync(stream.fileno())
                    try:
                        os.link(temporary_path, item.path)
                    except FileExistsError:
                        verify_existing(item)
                        continue
                    created.append(item.path)
                finally:
                    try:
                        temporary_path.unlink()
                    except OSError:
                        pass
            return created
        except Exception:
            self.cleanup_created(created)
            raise

    @staticmethod
    def cleanup_created(created: Iterable[Path]) -> None:
        for path in created:
            try:
                path.unlink()
            except OSError:
                pass


def parse_args(argv: list[str]) -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="PowerPoint（.pptx/.pptm）を解析し、Markdown変換用JSONを出力します。"
    )
    parser.add_argument("presentation", help="解析するPowerPointファイル")
    parser.add_argument(
        "--output",
        "-o",
        help="JSONの出力先。省略時は標準出力へ出力します。",
    )
    parser.add_argument(
        "--force",
        action="store_true",
        help="既存のJSON出力ファイルを原子的に置き換えます。",
    )
    parser.add_argument(
        "--slide",
        action="append",
        type=int,
        default=[],
        help="抽出する1始まりのスライド番号。複数回指定できます。",
    )
    parser.add_argument(
        "--include-hidden-slides",
        action="store_true",
        help="非表示スライドも抽出対象に含めます。",
    )
    parser.add_argument(
        "--include-notes",
        action="store_true",
        help="対象スライドの発表者ノート本文も抽出します。",
    )
    parser.add_argument(
        "--include-hidden-shapes",
        action="store_true",
        help="対象スライドの非表示図形も抽出します。",
    )
    parser.add_argument(
        "--include-off-slide-shapes",
        action="store_true",
        help="スライド領域外に完全に置かれた図形も抽出します。",
    )
    parser.add_argument(
        "--include-document-properties",
        action="store_true",
        help="作成者、タイトル、更新日時などの文書プロパティを抽出します。",
    )
    parser.add_argument(
        "--extract-images",
        metavar="DIRECTORY",
        help="対象スライドの埋め込み画像を指定ディレクトリへ抽出します。",
    )
    parser.add_argument(
        "--max-file-mb",
        type=int,
        default=DEFAULT_MAX_FILE_MB,
        help=f"入力ファイルの上限MiB（既定: {DEFAULT_MAX_FILE_MB}）。",
    )
    parser.add_argument(
        "--max-uncompressed-mb",
        type=int,
        default=DEFAULT_MAX_UNCOMPRESSED_MB,
        help=f"ZIP展開後の合計サイズ上限MiB（既定: {DEFAULT_MAX_UNCOMPRESSED_MB}）。",
    )
    parser.add_argument(
        "--max-slides",
        type=int,
        default=DEFAULT_MAX_SLIDES,
        help=f"全スライド数上限（既定: {DEFAULT_MAX_SLIDES}）。",
    )
    parser.add_argument(
        "--max-shapes",
        type=int,
        default=DEFAULT_MAX_SHAPES,
        help=f"抽出対象の図形数上限（既定: {DEFAULT_MAX_SHAPES}）。",
    )
    parser.add_argument(
        "--max-table-cells",
        type=int,
        default=DEFAULT_MAX_TABLE_CELLS,
        help=f"抽出対象の表セル数上限（既定: {DEFAULT_MAX_TABLE_CELLS}）。",
    )
    parser.add_argument(
        "--max-chart-points",
        type=int,
        default=DEFAULT_MAX_CHART_POINTS,
        help=f"抽出対象のグラフ点数上限（既定: {DEFAULT_MAX_CHART_POINTS}）。",
    )
    parser.add_argument(
        "--max-text-chars",
        type=int,
        default=DEFAULT_MAX_TEXT_CHARS,
        help=f"抽出対象の文字数上限（既定: {DEFAULT_MAX_TEXT_CHARS}）。",
    )
    parser.add_argument(
        "--max-extracted-images-mb",
        type=int,
        default=DEFAULT_MAX_EXTRACTED_IMAGES_MB,
        help=f"抽出画像の合計上限MiB（既定: {DEFAULT_MAX_EXTRACTED_IMAGES_MB}）。",
    )
    return parser.parse_args(argv)


def validate_args(args: argparse.Namespace) -> Path:
    path = Path(args.presentation)
    if not path.exists():
        raise InspectionError(f"入力ファイルが存在しません: {path}")
    if path.is_symlink():
        raise InspectionError("シンボリックリンクのPowerPointファイルは対象外です。")
    if not path.is_file():
        raise InspectionError(f"通常ファイルではありません: {path}")
    suffix = path.suffix.casefold()
    if suffix not in SUPPORTED_SUFFIXES:
        if suffix in LEGACY_SUFFIXES:
            raise InspectionError(
                f"{suffix}は対象外です。.pptxへ変換してから実行してください。"
            )
        if suffix in OTHER_OPENXML_SUFFIXES:
            raise InspectionError(
                f"{suffix}は標準対象外です。.pptxへ変換してから実行してください。"
            )
        raise InspectionError("対応形式は.pptxと.pptmです。")

    positive_options = (
        "max_file_mb",
        "max_uncompressed_mb",
        "max_slides",
        "max_shapes",
        "max_table_cells",
        "max_chart_points",
        "max_text_chars",
        "max_extracted_images_mb",
    )
    for option in positive_options:
        if getattr(args, option) <= 0:
            raise InspectionError(f"--{option.replace('_', '-')}には1以上を指定してください。")
    for slide_number in args.slide:
        if slide_number <= 0:
            raise InspectionError("--slideには1以上のスライド番号を指定してください。")

    size_limit = args.max_file_mb * 1024 * 1024
    if path.stat().st_size > size_limit:
        raise InspectionError(
            f"入力ファイルが上限{args.max_file_mb}MiBを超えています。"
            "安全性を確認してから上限を変更してください。"
        )
    if args.force and not args.output:
        raise InspectionError("--forceは--outputと同時に指定してください。")
    if args.output:
        output_path = Path(args.output)
        if output_path.resolve(strict=False) == path.resolve():
            raise InspectionError("出力先に入力PowerPointと同じパスは指定できません。")
        if output_path.exists():
            try:
                if output_path.samefile(path):
                    raise InspectionError(
                        "出力先に入力PowerPointと同じファイルは指定できません。"
                    )
            except OSError:
                pass
        if output_path.suffix.casefold() != ".json":
            raise InspectionError("--outputには拡張子.jsonのファイルを指定してください。")
        if output_path.exists() and output_path.is_dir():
            raise InspectionError("JSONの出力先にディレクトリは指定できません。")
        if os.path.lexists(output_path) and not args.force:
            raise InspectionError(
                "JSONの出力先が既に存在します。上書きする場合だけ--forceを指定してください。"
            )
    if args.extract_images:
        image_dir = Path(args.extract_images)
        if image_dir.exists() and image_dir.is_symlink():
            raise InspectionError("画像の出力先にシンボリックリンクは指定できません。")
        if image_dir.exists() and not image_dir.is_dir():
            raise InspectionError("画像の出力先はディレクトリである必要があります。")
    return path


def local_name(tag: str) -> str:
    return tag.rsplit("}", 1)[-1] if "}" in tag else tag.rsplit(":", 1)[-1]


def escape_untrusted_text(value: Any) -> str:
    text = str(value).replace("\r\n", "\n").replace("\r", "\n").replace("\v", "\n")
    result: list[str] = []
    for character in text:
        code = ord(character)
        if code in BIDI_AND_DIRECTIONAL_CONTROLS:
            result.append(f"\\u{code:04X}")
        elif code < 32 and character not in {"\n", "\t"}:
            result.append(f"\\u{code:04X}")
        elif code == 127:
            result.append("\\u007F")
        else:
            result.append(character)
    return "".join(result)


def relationship_type_label(value: str) -> str:
    candidate = value.rsplit("/", 1)[-1]
    if candidate in SAFE_RELATIONSHIP_TYPE_LABELS:
        return candidate
    return "other"


def add_text(value: Any, limits: Limits) -> str:
    text = escape_untrusted_text(value)
    limits.add_text(text)
    return text


def enum_name(value: Any) -> str | None:
    if value is None:
        return None
    name = getattr(value, "name", None)
    if isinstance(name, str):
        return name
    text = str(value)
    match = re.match(r"^([A-Z][A-Z0-9_]*)", text)
    return match.group(1) if match else escape_untrusted_text(text)


def scalar_json(value: Any) -> Any:
    if value is None or isinstance(value, (bool, int, str)):
        return value
    if isinstance(value, float):
        if math.isnan(value):
            return "NaN"
        if math.isinf(value):
            return "Infinity" if value > 0 else "-Infinity"
        return value
    if isinstance(value, (datetime, date, time)):
        return value.isoformat()
    if isinstance(value, (list, tuple)):
        return [scalar_json(item) for item in value]
    return escape_untrusted_text(value)


def xml_scan(member_stream: Any, max_bytes: int) -> tuple[bool, bool, int]:
    """XMLらしさとDTD/ENTITY宣言の有無をストリーミング判定する。"""

    overlap = b""
    first_chunk = True
    looks_like_xml: bool | None = None
    bytes_read = 0
    while True:
        chunk = member_stream.read(64 * 1024)
        if not chunk:
            return bool(looks_like_xml), False, bytes_read
        bytes_read += len(chunk)
        if bytes_read > max_bytes:
            raise InspectionError("ZIP展開後の合計サイズが上限を超えています。")
        if first_chunk:
            first_chunk = False
            if chunk.startswith(b"\x4c\x6f\xa7\x94"):
                return True, True, bytes_read
        data = overlap + chunk
        normalized = data.replace(b"\x00", b"").upper()
        if looks_like_xml is None:
            prefix = normalized
            for marker in (b"\xef\xbb\xbf", b"\xff\xfe", b"\xfe\xff"):
                if prefix.startswith(marker.upper()):
                    prefix = prefix[len(marker) :]
                    break
            significant = prefix.lstrip(b" \t\r\n")
            if significant:
                looks_like_xml = significant.startswith(b"<")
        if looks_like_xml and (b"<!DOCTYPE" in normalized or b"<!ENTITY" in normalized):
            return True, True, bytes_read
        overlap = data[-64:]


def fully_unquote(value: str) -> str:
    decoded = value
    for _ in range(4):
        next_value = unquote(decoded)
        if next_value == decoded:
            return decoded
        decoded = next_value
    return decoded


def canonical_member_key(name: str) -> str:
    decoded = fully_unquote(name)
    if "\x00" in decoded or "\\" in decoded:
        raise InspectionError(f"ZIP部品名のエンコード後パスが不正です: {name}")
    if decoded.startswith("/") or re.match(r"^[A-Za-z]:", decoded):
        raise InspectionError(f"ZIP部品名のエンコード後パスが絶対パスです: {name}")
    parts = PurePosixPath(decoded).parts
    if any(part in {"", ".", ".."} for part in parts):
        raise InspectionError(f"ZIP部品名のエンコード後パスに不正な要素があります: {name}")
    normalized = posixpath.normpath(decoded)
    if normalized != decoded or normalized.startswith("../") or normalized == "..":
        raise InspectionError(f"ZIP部品名のエンコード後パスがパッケージ外を参照します: {name}")
    return unicodedata.normalize("NFC", normalized).casefold()


def validate_member_name(name: str) -> str:
    if not name or "\x00" in name or "\\" in name:
        raise InspectionError("ZIP部品名に不正な文字が含まれています。")
    if name.startswith("/") or re.match(r"^[A-Za-z]:", name):
        raise InspectionError(f"ZIP部品名に絶対パスは使用できません: {name}")
    stripped = name[:-1] if name.endswith("/") else name
    parts = PurePosixPath(stripped).parts
    if any(part in {"", ".", ".."} for part in parts):
        raise InspectionError(f"ZIP部品名に不正なパス要素があります: {name}")
    normalized = posixpath.normpath(stripped)
    if normalized != stripped or normalized.startswith("../") or normalized == "..":
        raise InspectionError(f"ZIP部品名がパッケージ外を参照します: {name}")
    return stripped


def relationship_source_directory(relationship_name: str) -> str:
    if relationship_name == "_rels/.rels":
        return ""
    path = PurePosixPath(relationship_name)
    if path.parent.name != "_rels" or not path.name.endswith(".rels"):
        raise InspectionError(f"Relationship部品の配置が不正です: {relationship_name}")
    source_name = path.name[: -len(".rels")]
    source_part = path.parent.parent / source_name
    return str(source_part.parent)


def resolve_internal_relationship(
    relationship_name: str,
    target: str,
    canonical_member_names: set[str],
) -> str:
    decoded = fully_unquote(target.split("#", 1)[0])
    if not decoded:
        raise InspectionError(f"空の内部Relationshipがあります: {relationship_name}")
    if "\x00" in decoded or "\\" in decoded:
        raise InspectionError(f"内部Relationshipに不正な文字があります: {relationship_name}")
    parsed = urlsplit(decoded)
    if (
        parsed.scheme
        or parsed.netloc
        or parsed.query
        or decoded.startswith("/")
        or re.match(r"^[A-Za-z]:", decoded)
    ):
        raise InspectionError(f"内部Relationshipに絶対パスまたはURIがあります: {relationship_name}")
    base = relationship_source_directory(relationship_name)
    normalized = posixpath.normpath(posixpath.join(base, decoded))
    if normalized in {"", ".", ".."} or normalized.startswith("../"):
        raise InspectionError(f"内部Relationshipがパッケージ外を参照します: {relationship_name}")
    canonical_target = unicodedata.normalize("NFC", normalized).casefold()
    if canonical_target not in canonical_member_names:
        raise InspectionError(
            f"内部Relationshipの参照先が存在しません: {relationship_name} -> {normalized}"
        )
    return normalized


class PackagePictureResolver:
    """Picture図形を検証済みOOXML部品へ対応付け、外部参照は開かない。"""

    def __init__(self, source_bytes: bytes):
        self._stream = io.BytesIO(source_bytes)
        self._archive = zipfile.ZipFile(self._stream)
        self._member_by_key = {
            canonical_member_key(info.filename): info.filename
            for info in self._archive.infolist()
            if not info.is_dir()
        }
        self._member_keys = set(self._member_by_key)
        self._relationship_cache: dict[
            str, dict[str, tuple[str, str, bool]]
        ] = {}
        self._default_content_types: dict[str, str] = {}
        self._override_content_types: dict[str, str] = {}
        self._load_content_types()

    def close(self) -> None:
        self._archive.close()
        self._stream.close()

    def __enter__(self) -> "PackagePictureResolver":
        return self

    def __exit__(self, exc_type: Any, exc: Any, traceback: Any) -> None:
        self.close()

    def _load_content_types(self) -> None:
        try:
            root = ElementTree.fromstring(self._archive.read("[Content_Types].xml"))
        except (KeyError, ElementTree.ParseError) as exc:
            raise InspectionError("[Content_Types].xmlを解析できません。") from exc
        for element in root:
            name = local_name(element.tag)
            content_type = element.attrib.get("ContentType", "").strip()
            if not content_type:
                continue
            if name == "Default":
                extension = element.attrib.get("Extension", "").lstrip(".").casefold()
                if extension:
                    self._default_content_types[extension] = content_type
            elif name == "Override":
                part_name = element.attrib.get("PartName", "").lstrip("/")
                if not part_name:
                    continue
                try:
                    key = canonical_member_key(part_name)
                except InspectionError:
                    continue
                self._override_content_types[key] = content_type

    @staticmethod
    def _relationship_part_name(source_part_name: str) -> str:
        source_path = PurePosixPath(source_part_name.lstrip("/"))
        return (
            source_path.parent / "_rels" / f"{source_path.name}.rels"
        ).as_posix()

    def _relationships(self, source_part_name: str) -> dict[str, tuple[str, str, bool]]:
        relationship_name = self._relationship_part_name(source_part_name)
        if relationship_name in self._relationship_cache:
            return self._relationship_cache[relationship_name]
        actual_name = self._member_by_key.get(canonical_member_key(relationship_name))
        if actual_name is None:
            self._relationship_cache[relationship_name] = {}
            return {}
        try:
            root = ElementTree.fromstring(self._archive.read(actual_name))
        except ElementTree.ParseError as exc:
            raise InspectionError(
                f"Relationship XMLを解析できません: {actual_name}"
            ) from exc
        relationships: dict[str, tuple[str, str, bool]] = {}
        for element in root:
            if local_name(element.tag) != "Relationship":
                continue
            relationship_id = element.attrib.get("Id", "")
            if not relationship_id:
                continue
            if relationship_id in relationships:
                raise InspectionError(
                    f"Relationship IDが重複しています: {actual_name}"
                )
            relationships[relationship_id] = (
                element.attrib.get("Type", ""),
                element.attrib.get("Target", ""),
                element.attrib.get("TargetMode", "").casefold() == "external",
            )
        self._relationship_cache[relationship_name] = relationships
        return relationships

    @staticmethod
    def _picture_relationship(shape: Any) -> tuple[str, str] | None:
        if local_name(shape.element.tag) != "pic":
            return None
        embedded: tuple[str, str] | None = None
        extended_embedded: tuple[str, str] | None = None
        linked: tuple[str, str] | None = None
        for element in shape.element.iter():
            element_name = local_name(element.tag)
            if element_name not in {"blip", "svgBlip"}:
                continue
            for attribute, relationship_id in element.attrib.items():
                attribute_name = local_name(attribute)
                if attribute_name == "embed" and relationship_id:
                    if element_name == "blip":
                        embedded = (relationship_id, "embed")
                    else:
                        extended_embedded = (relationship_id, "extended_embed")
                if attribute_name == "link" and relationship_id:
                    linked = (relationship_id, "link")
        return extended_embedded or embedded or linked

    def resolve(self, shape: Any, source_part_name: str) -> PackageImage | None:
        relationship_reference = self._picture_relationship(shape)
        if relationship_reference is None:
            if local_name(shape.element.tag) != "pic":
                return None
            return PackageImage(
                None,
                None,
                None,
                None,
                "Picture図形に画像Relationship IDがありません。",
            )
        relationship_id, reference_kind = relationship_reference
        relationships = self._relationships(source_part_name)
        relationship = relationships.get(relationship_id)
        if relationship is None:
            return PackageImage(
                None,
                None,
                None,
                None,
                "Picture図形の画像Relationshipを解決できません。",
            )
        relationship_type, target, external = relationship
        if external:
            return PackageImage(
                None,
                None,
                None,
                None,
                "外部画像Relationshipのリンク先にはアクセスしていません。",
            )
        if relationship_type not in IMAGE_RELATIONSHIP_TYPES:
            return PackageImage(
                None,
                None,
                None,
                None,
                "Picture図形のRelationship Typeが埋め込み画像ではないため開いていません。",
            )
        relationship_name = self._relationship_part_name(source_part_name)
        resolved_name = resolve_internal_relationship(
            relationship_name,
            target,
            self._member_keys,
        )
        actual_name = self._member_by_key[canonical_member_key(resolved_name)]
        data = self._archive.read(actual_name)
        filename = PurePosixPath(actual_name).name
        suffix = PurePosixPath(filename).suffix.lstrip(".").casefold() or None
        member_key = canonical_member_key(actual_name)
        content_type = self._override_content_types.get(member_key)
        if content_type is None and suffix is not None:
            content_type = self._default_content_types.get(suffix)
        warning = None
        if reference_kind == "link":
            warning = "内部部品へのリンク画像として記録しました。"
        elif reference_kind == "extended_embed":
            warning = (
                "Picture図形の拡張画像RelationshipをPNGフォールバックより優先して"
                "記録しました。"
            )
        return PackageImage(
            filename,
            content_type,
            suffix,
            data,
            warning,
            reference_kind != "extended_embed",
        )


def count_slide_xml_features(data: bytes) -> tuple[bool, bool]:
    try:
        root = ElementTree.fromstring(data)
    except ElementTree.ParseError as exc:
        raise InspectionError("スライドXMLを解析できません。") from exc
    has_timing = any(local_name(element.tag) == "timing" for element in root.iter())
    has_transition = any(local_name(element.tag) == "transition" for element in root.iter())
    return has_timing, has_transition


def inspect_archive(
    source_bytes: bytes,
    max_uncompressed_mb: int,
    *,
    macro_capable_suffix: bool,
) -> dict[str, Any]:
    try:
        archive = zipfile.ZipFile(io.BytesIO(source_bytes))
    except zipfile.BadZipFile as exc:
        raise InspectionError("PowerPointファイルが有効なZIP形式ではありません。") from exc

    with archive:
        infos = archive.infolist()
        if len(infos) > MAX_ARCHIVE_PARTS:
            raise InspectionError(
                f"ZIP部品数が上限{MAX_ARCHIVE_PARTS}を超えています。"
            )
        member_names: set[str] = set()
        canonical_member_names: set[str] = set()
        total_uncompressed = 0
        total_compressed = 0

        for info in infos:
            name = validate_member_name(info.filename)
            canonical_key = canonical_member_key(name)
            if name in member_names or canonical_key in canonical_member_names:
                raise InspectionError(f"ZIP部品名が重複または大小文字衝突しています: {name}")
            member_names.add(name)
            canonical_member_names.add(canonical_key)
            if info.flag_bits & 0x1:
                raise InspectionError("暗号化されたZIP部品は解析できません。")
            unix_mode = (info.external_attr >> 16) & 0xFFFF
            if unix_mode and stat.S_IFMT(unix_mode) == stat.S_IFLNK:
                raise InspectionError(f"ZIP内のシンボリックリンクは対象外です: {name}")
            total_uncompressed += info.file_size
            total_compressed += info.compress_size
            if info.file_size > 1024 * 1024:
                if info.compress_size == 0:
                    raise InspectionError(f"異常な圧縮部品を検出しました: {name}")
                ratio = info.file_size / info.compress_size
                if ratio > MAX_SUSPICIOUS_COMPRESSION_RATIO:
                    raise InspectionError(f"異常な圧縮率の部品を検出しました: {name}")

        uncompressed_limit = max_uncompressed_mb * 1024 * 1024
        if total_uncompressed > uncompressed_limit:
            raise InspectionError(
                f"ZIP展開後の合計サイズが上限{max_uncompressed_mb}MiBを超えています。"
            )
        required = {"[Content_Types].xml", "_rels/.rels", "ppt/presentation.xml"}
        missing = sorted(required - member_names)
        if missing:
            raise InspectionError(
                "PowerPointに必要な部品がありません: " + ", ".join(missing)
            )

        actual_uncompressed = 0
        for info in infos:
            if info.is_dir():
                continue
            with archive.open(info) as stream:
                looks_xml, forbidden, member_bytes = xml_scan(
                    stream, uncompressed_limit - actual_uncompressed
                )
            actual_uncompressed += member_bytes
            if looks_xml and info.file_size > MAX_XML_PART_BYTES:
                raise InspectionError(f"XML部品が上限を超えています: {info.filename}")
            if forbidden:
                raise InspectionError(
                    f"DTDまたはXMLエンティティ宣言を検出しました: {info.filename}"
                )

        external_relationship_count = 0
        internal_relationship_count = 0
        relationship_types: dict[str, int] = {}
        has_vba_relationship = False
        for info in infos:
            if info.is_dir() or not info.filename.casefold().endswith(".rels"):
                continue
            if info.file_size > MAX_RELATIONSHIPS_BYTES:
                raise InspectionError(f"Relationship部品が上限を超えています: {info.filename}")
            data = archive.read(info)
            try:
                root = ElementTree.fromstring(data)
            except ElementTree.ParseError as exc:
                raise InspectionError(f"Relationship XMLを解析できません: {info.filename}") from exc
            for element in root.iter():
                if local_name(element.tag) != "Relationship":
                    continue
                target = element.attrib.get("Target", "")
                relationship_type = relationship_type_label(element.attrib.get("Type", ""))
                has_vba_relationship = has_vba_relationship or relationship_type == "vbaProject"
                relationship_types[relationship_type] = (
                    relationship_types.get(relationship_type, 0) + 1
                )
                if element.attrib.get("TargetMode", "").casefold() == "external":
                    external_relationship_count += 1
                    continue
                resolve_internal_relationship(
                    info.filename, target, canonical_member_names
                )
                internal_relationship_count += 1

        lower_names = [name.casefold() for name in member_names]
        try:
            content_types_root = ElementTree.fromstring(archive.read("[Content_Types].xml"))
        except ElementTree.ParseError as exc:
            raise InspectionError("[Content_Types].xmlを解析できません。") from exc
        has_vba_content_type = any(
            "vbaproject" in element.attrib.get("ContentType", "").casefold()
            or "macroenabled" in element.attrib.get("ContentType", "").casefold()
            for element in content_types_root.iter()
        )
        slide_xml_names = sorted(
            name
            for name in member_names
            if re.match(r"^ppt/slides/slide\d+\.xml$", name, re.IGNORECASE)
        )
        slides_with_timing = 0
        slides_with_transition = 0
        for name in slide_xml_names:
            timing, transition = count_slide_xml_features(archive.read(name))
            slides_with_timing += int(timing)
            slides_with_transition += int(transition)

        media_names = [name for name in member_names if name.casefold().startswith("ppt/media/")]
        image_parts = sum(Path(name).suffix.casefold() in IMAGE_SUFFIXES for name in media_names)
        audio_video_parts = sum(
            Path(name).suffix.casefold() in AUDIO_VIDEO_SUFFIXES for name in media_names
        )

        return {
            "part_count": len(infos),
            "compressed_bytes": total_compressed,
            "uncompressed_bytes": actual_uncompressed,
            "declared_uncompressed_bytes": total_uncompressed,
            "internal_relationship_count": internal_relationship_count,
            "external_relationship_count": external_relationship_count,
            "relationship_types": dict(sorted(relationship_types.items())),
            "slide_parts": len(slide_xml_names),
            "notes_slide_parts": sum(
                bool(re.match(r"^ppt/notesslides/notesslide\d+\.xml$", name, re.IGNORECASE))
                for name in member_names
            ),
            "chart_parts": sum(name.startswith("ppt/charts/") for name in lower_names),
            "diagram_parts": sum(name.startswith("ppt/diagrams/") for name in lower_names),
            "embedding_parts": sum(name.startswith("ppt/embeddings/") for name in lower_names),
            "active_x_parts": sum(name.startswith("ppt/activex/") for name in lower_names),
            "comment_parts": sum(
                name.startswith("ppt/comments/")
                or name.startswith("ppt/threadedcomments/")
                or name == "ppt/commentauthors.xml"
                or name.startswith("ppt/persons/")
                for name in lower_names
            ),
            "custom_xml_parts": sum(name.startswith("customxml/") for name in lower_names),
            "image_parts": image_parts,
            "audio_video_parts": audio_video_parts,
            "other_media_parts": max(len(media_names) - image_parts - audio_video_parts, 0),
            "has_vba": (
                macro_capable_suffix
                or has_vba_relationship
                or has_vba_content_type
                or any(name.endswith("vbaproject.bin") for name in lower_names)
            ),
            "slides_with_timing": slides_with_timing,
            "slides_with_transition": slides_with_transition,
        }


def parse_on_off(value: Any, default: bool = True) -> bool:
    if value is None:
        return default
    return str(value).strip().casefold() not in {"0", "false", "off", "no"}


def slide_is_hidden(slide: Any) -> bool:
    return not parse_on_off(slide.element.get("show"), default=True)


def shape_cnvpr(shape: Any) -> Any | None:
    for element in shape.element.iter():
        if local_name(element.tag) == "cNvPr":
            return element
    return None


def shape_is_hidden(shape: Any) -> bool:
    cnvpr = shape_cnvpr(shape)
    return False if cnvpr is None else parse_on_off(cnvpr.get("hidden"), default=False)


def shape_alt_text(shape: Any, limits: Limits) -> dict[str, str] | None:
    cnvpr = shape_cnvpr(shape)
    if cnvpr is None:
        return None
    result: dict[str, str] = {}
    description = cnvpr.get("descr")
    title = cnvpr.get("title")
    if description:
        result["description"] = add_text(description, limits)
    if title:
        result["title"] = add_text(title, limits)
    return result or None


def length_dict(value: Any) -> dict[str, Any] | None:
    if value is None:
        return None
    integer = int(value)
    return {"emu": integer, "inches": round(integer / EMU_PER_INCH, 6)}


def shape_corners(
    shape: Any,
    parent_transform: AffineTransform | None,
) -> list[tuple[float, float]] | None:
    if parent_transform is None:
        return None
    values = [getattr(shape, name, None) for name in ("left", "top", "width", "height")]
    if any(value is None for value in values):
        return None
    left, top, width, height = map(float, values)
    corners = [
        (left, top),
        (left + width, top),
        (left + width, top + height),
        (left, top + height),
    ]
    rotation = float(getattr(shape, "rotation", 0.0) or 0.0)
    if rotation:
        radians = math.radians(rotation)
        cosine = math.cos(radians)
        sine = math.sin(radians)
        center_x = left + width / 2
        center_y = top + height / 2
        rotated: list[tuple[float, float]] = []
        for x, y in corners:
            relative_x = x - center_x
            relative_y = y - center_y
            rotated.append(
                (
                    center_x + cosine * relative_x - sine * relative_y,
                    center_y + sine * relative_x + cosine * relative_y,
                )
            )
        corners = rotated
    return [parent_transform.point(x, y) for x, y in corners]


def bounds_from_corners(
    corners: list[tuple[float, float]],
) -> tuple[float, float, float, float]:
    x_values = [point[0] for point in corners]
    y_values = [point[1] for point in corners]
    return min(x_values), min(y_values), max(x_values), max(y_values)


def direct_xml_child(element: Any, name: str) -> Any | None:
    for child in element:
        if local_name(child.tag) == name:
            return child
    return None


def group_child_transform(
    shape: Any,
    parent_transform: AffineTransform | None,
) -> AffineTransform | None:
    if parent_transform is None:
        return None
    group_properties = direct_xml_child(shape.element, "grpSpPr")
    xfrm = None if group_properties is None else direct_xml_child(group_properties, "xfrm")
    if xfrm is None:
        return None
    child_offset = direct_xml_child(xfrm, "chOff")
    child_extent = direct_xml_child(xfrm, "chExt")
    values = [getattr(shape, name, None) for name in ("left", "top", "width", "height")]
    if child_offset is None or child_extent is None or any(value is None for value in values):
        return None
    try:
        child_x = float(child_offset.get("x"))
        child_y = float(child_offset.get("y"))
        child_width = float(child_extent.get("cx"))
        child_height = float(child_extent.get("cy"))
        left, top, width, height = map(float, values)
    except (TypeError, ValueError):
        return None
    if child_width == 0 or child_height == 0:
        return None

    scale_x = width / child_width
    scale_y = height / child_height
    mapping = AffineTransform(
        a=scale_x,
        d=scale_y,
        e=left - child_x * scale_x,
        f=top - child_y * scale_y,
    )
    center_x = left + width / 2
    center_y = top + height / 2
    if parse_on_off(xfrm.get("flipH"), default=False):
        mapping = AffineTransform(a=-1.0, d=1.0, e=2 * center_x).compose(mapping)
    if parse_on_off(xfrm.get("flipV"), default=False):
        mapping = AffineTransform(a=1.0, d=-1.0, f=2 * center_y).compose(mapping)

    rotation = float(getattr(shape, "rotation", 0.0) or 0.0)
    if rotation:
        radians = math.radians(rotation)
        cosine = math.cos(radians)
        sine = math.sin(radians)
        rotate = AffineTransform(
            a=cosine,
            b=sine,
            c=-sine,
            d=cosine,
            e=center_x - cosine * center_x + sine * center_y,
            f=center_y - sine * center_x - cosine * center_y,
        )
        mapping = rotate.compose(mapping)
    return parent_transform.compose(mapping)


def position_dict(
    bounds: tuple[float, float, float, float],
    slide_width: int,
    slide_height: int,
    rotation: Any,
) -> dict[str, Any]:
    left_value, top_value, right_value, bottom_value = bounds
    left = int(round(left_value))
    top = int(round(top_value))
    width = int(round(right_value - left_value))
    height = int(round(bottom_value - top_value))
    result: dict[str, Any] = {
        "left": length_dict(left),
        "top": length_dict(top),
        "width": length_dict(width),
        "height": length_dict(height),
    }
    result["relative"] = {
        "left": round(left / slide_width, 6) if slide_width else None,
        "top": round(top / slide_height, 6) if slide_height else None,
        "width": round(width / slide_width, 6) if slide_width else None,
        "height": round(height / slide_height, 6) if slide_height else None,
    }
    if rotation is not None:
        result["rotation_degrees"] = float(rotation)
    return result


def off_slide_state(
    corners: list[tuple[float, float]] | None,
    slide_width: int,
    slide_height: int,
) -> tuple[bool, bool]:
    if corners is None:
        return False, False
    slide_corners = [
        (0.0, 0.0),
        (float(slide_width), 0.0),
        (float(slide_width), float(slide_height)),
        (0.0, float(slide_height)),
    ]

    def polygon_area(polygon: list[tuple[float, float]]) -> float:
        return abs(
            sum(
                x1 * y2 - x2 * y1
                for (x1, y1), (x2, y2) in zip(
                    polygon, polygon[1:] + polygon[:1]
                )
            )
        ) / 2

    def overlaps_with_positive_area(
        first: list[tuple[float, float]],
        second: list[tuple[float, float]],
    ) -> bool:
        if polygon_area(first) < 0.5:
            endpoints = max(
                (
                    (start, end)
                    for start in first
                    for end in first
                ),
                key=lambda pair: (pair[1][0] - pair[0][0]) ** 2
                + (pair[1][1] - pair[0][1]) ** 2,
            )
            (start_x, start_y), (end_x, end_y) = endpoints
            delta_x = end_x - start_x
            delta_y = end_y - start_y
            if delta_x == 0 and delta_y == 0:
                return (
                    0 <= start_x <= slide_width
                    and 0 <= start_y <= slide_height
                )
            minimum_t = 0.0
            maximum_t = 1.0
            for coefficient, distance in (
                (-delta_x, start_x),
                (delta_x, slide_width - start_x),
                (-delta_y, start_y),
                (delta_y, slide_height - start_y),
            ):
                if coefficient == 0:
                    if distance < 0:
                        return False
                    continue
                ratio = distance / coefficient
                if coefficient < 0:
                    minimum_t = max(minimum_t, ratio)
                else:
                    maximum_t = min(maximum_t, ratio)
                if minimum_t > maximum_t:
                    return False
            return True
        if polygon_area(second) < 0.5:
            return False
        for polygon in (first, second):
            for start, end in zip(polygon, polygon[1:] + polygon[:1]):
                axis_x = -(end[1] - start[1])
                axis_y = end[0] - start[0]
                if axis_x == 0 and axis_y == 0:
                    continue
                first_projection = [x * axis_x + y * axis_y for x, y in first]
                second_projection = [x * axis_x + y * axis_y for x, y in second]
                if max(first_projection) <= min(second_projection) or max(
                    second_projection
                ) <= min(first_projection):
                    return False
        return True

    fully = not overlaps_with_positive_area(corners, slide_corners)
    inside = all(
        0 <= x <= slide_width and 0 <= y <= slide_height for x, y in corners
    )
    partially = not fully and not inside
    return fully, partially


def bullet_info(paragraph: Any) -> dict[str, Any]:
    p_pr = getattr(paragraph._p, "pPr", None)
    if p_pr is None:
        return {"type": "inherited_or_unspecified"}
    for child in p_pr:
        name = local_name(child.tag)
        if name == "buNone":
            return {"type": "none"}
        if name == "buChar":
            return {"type": "character", "character": escape_untrusted_text(child.get("char", ""))}
        if name == "buAutoNum":
            return {"type": "numbered", "numbering": escape_untrusted_text(child.get("type", ""))}
    return {"type": "inherited_or_unspecified"}


def run_dict(run: Any, limits: Limits) -> dict[str, Any]:
    result: dict[str, Any] = {"text": escape_untrusted_text(run.text)}
    font = run.font
    style: dict[str, Any] = {}
    if font.bold is not None:
        style["bold"] = bool(font.bold)
    if font.italic is not None:
        style["italic"] = bool(font.italic)
    if font.underline is not None:
        style["underline"] = scalar_json(font.underline)
    if font.name:
        style["font_name"] = add_text(font.name, limits)
    if font.size is not None:
        style["font_size_pt"] = round(float(font.size.pt), 3)
    if style:
        result["style"] = style
    try:
        address = run.hyperlink.address
    except Exception:
        address = None
    if address:
        result["hyperlink"] = add_text(address, limits)
    return result


def text_frame_dict(text_frame: Any, limits: Limits) -> dict[str, Any]:
    paragraphs: list[dict[str, Any]] = []
    full_text_parts: list[str] = []
    for index, paragraph in enumerate(text_frame.paragraphs):
        paragraph_text = add_text(paragraph.text, limits)
        full_text_parts.append(paragraph_text)
        item: dict[str, Any] = {
            "index": index,
            "text": paragraph_text,
            "level": int(paragraph.level),
            "bullet": bullet_info(paragraph),
            "runs": [run_dict(run, limits) for run in paragraph.runs],
        }
        alignment = enum_name(paragraph.alignment)
        if alignment is not None:
            item["alignment"] = alignment
        paragraphs.append(item)
    result: dict[str, Any] = {
        "text": "\n".join(full_text_parts),
        "paragraphs": paragraphs,
    }
    vertical_anchor = enum_name(getattr(text_frame, "vertical_anchor", None))
    if vertical_anchor is not None:
        result["vertical_anchor"] = vertical_anchor
    if getattr(text_frame, "word_wrap", None) is not None:
        result["word_wrap"] = bool(text_frame.word_wrap)
    return result


def table_dict(table: Any, limits: Limits) -> dict[str, Any]:
    row_count = len(table.rows)
    column_count = len(table.columns)
    limits.add_table_cells(row_count * column_count)
    rows: list[list[dict[str, Any]]] = []
    for row_index in range(row_count):
        row_values: list[dict[str, Any]] = []
        for column_index in range(column_count):
            cell = table.cell(row_index, column_index)
            item: dict[str, Any] = {
                "row": row_index + 1,
                "column": column_index + 1,
                "text": text_frame_dict(cell.text_frame, limits),
                "is_spanned": bool(cell.is_spanned),
                "is_merge_origin": bool(cell.is_merge_origin),
            }
            if cell.is_merge_origin:
                item["span_height"] = int(cell.span_height)
                item["span_width"] = int(cell.span_width)
            row_values.append(item)
        rows.append(row_values)
    return {
        "row_count": row_count,
        "column_count": column_count,
        "row_heights": [length_dict(row.height) for row in table.rows],
        "column_widths": [length_dict(column.width) for column in table.columns],
        "rows": rows,
    }


def safe_chart_values(values: Iterable[Any], limits: Limits) -> list[Any]:
    result = [chart_json_value(value, limits) for value in values]
    limits.add_chart_points(len(result))
    return result


def chart_json_value(value: Any, limits: Limits) -> Any:
    if isinstance(value, str):
        return add_text(value, limits)
    if isinstance(value, (list, tuple)):
        return [chart_json_value(item, limits) for item in value]
    return scalar_json(value)


def chart_dict(chart: Any, limits: Limits) -> dict[str, Any]:
    warnings: list[str] = []
    result: dict[str, Any] = {
        "chart_type": None,
        "title": None,
        "series": [],
        "categories": [],
        "cached_values_only": True,
        "warnings": warnings,
    }
    try:
        result["chart_type"] = enum_name(chart.chart_type)
    except InspectionError:
        raise
    except Exception as exc:
        warnings.append(f"グラフ種類を取得できません: {type(exc).__name__}")
    try:
        if chart.has_title:
            title_frame = chart.chart_title.text_frame
            result["title"] = text_frame_dict(title_frame, limits)
    except InspectionError:
        raise
    except Exception as exc:
        warnings.append(f"グラフタイトルを取得できません: {type(exc).__name__}")

    try:
        series_items: list[dict[str, Any]] = []
        for index, series in enumerate(chart.series):
            item: dict[str, Any] = {"index": index}
            try:
                item["name"] = add_text(series.name, limits)
            except InspectionError:
                raise
            except Exception:
                item["name"] = None
            try:
                item["values"] = safe_chart_values(series.values, limits)
            except InspectionError:
                raise
            except Exception as exc:
                item["values"] = None
                item["warning"] = f"系列値を取得できません: {type(exc).__name__}"
            series_items.append(item)
        result["series"] = series_items
    except InspectionError:
        raise
    except Exception as exc:
        warnings.append(f"グラフ系列を取得できません: {type(exc).__name__}")

    try:
        plots: list[dict[str, Any]] = []
        for index, plot in enumerate(chart.plots):
            item: dict[str, Any] = {"index": index}
            categories = getattr(plot, "categories", None)
            if categories is not None:
                labels = [
                    chart_json_value(label, limits) for label in categories.flattened_labels
                ]
                limits.add_chart_points(len(labels))
                item["labels"] = labels
            plots.append(item)
        result["categories"] = plots
    except InspectionError:
        raise
    except Exception as exc:
        warnings.append(f"グラフカテゴリを取得できません: {type(exc).__name__}")
    warnings.append("グラフ値は保存済みキャッシュであり、リンク元を再計算していません。")
    warnings.append("軸、単位、凡例、強調表現はスライドを視覚確認してください。")
    return result


def click_action_dict(shape: Any, limits: Limits) -> dict[str, Any] | None:
    try:
        click_action = shape.click_action
        action_name = enum_name(click_action.action)
    except Exception:
        return None
    if action_name in {None, "NONE"}:
        return None
    result: dict[str, Any] = {"action": action_name}
    try:
        address = click_action.hyperlink.address
    except Exception:
        address = None
    if address:
        result["hyperlink"] = add_text(address, limits)
    try:
        target_slide = click_action.target_slide
    except Exception:
        target_slide = None
    if target_slide is not None:
        result["target_slide_id"] = int(target_slide.slide_id)
    if action_name in {"RUN_MACRO", "RUN_PROGRAM", "OLE_VERB"}:
        result["unsafe_action"] = True
        result["warning"] = "アクションの対象は取得または実行していません。"
    return result


def picture_dict(
    shape: Any,
    limits: Limits,
    image_collector: ImageCollector,
    picture_resolver: PackagePictureResolver,
    source_part_name: str,
    slide_number: int,
) -> dict[str, Any] | None:
    package_image = picture_resolver.resolve(shape, source_part_name)
    if package_image is None:
        return None
    data = package_image.data
    extension = package_image.extension
    content_type = package_image.content_type
    extracted_path: str | None = None
    extractable = False
    not_extracted_reason = package_image.resolution_warning
    if data is not None:
        extractable, classification_reason = image_extraction_decision(
            extension, content_type, data
        )
        extracted_path, extraction_reason = image_collector.add(
            slide_number=slide_number,
            shape_id=int(shape.shape_id),
            extension=extension or "",
            content_type=content_type,
            data=data,
        )
        not_extracted_reason = extraction_reason or classification_reason

    pixel_size = None
    python_pptx_warning = None
    if data is not None and package_image.decode_with_python_pptx:
        try:
            image = shape.image
            pixel_size = {"width": int(image.size[0]), "height": int(image.size[1])}
        except Exception as exc:
            python_pptx_warning = (
                "python-pptxで画像寸法を取得できなかったため、"
                f"OOXML部品のメタデータだけを記録しました: {type(exc).__name__}"
            )
    elif data is not None:
        python_pptx_warning = (
            "python-pptxのshape.imageはフォールバック画像を示すため、"
            "拡張画像部品の画素寸法として使用していません。"
        )

    result: dict[str, Any] = {
        "filename": (
            add_text(package_image.filename, limits)
            if package_image.filename is not None
            else None
        ),
        "content_type": (
            escape_untrusted_text(content_type) if content_type is not None else None
        ),
        "extension": (
            escape_untrusted_text(extension) if extension is not None else None
        ),
        "bytes": len(data) if data is not None else None,
        "sha256": hashlib.sha256(data).hexdigest() if data is not None else None,
        "pixel_size": pixel_size,
        "extractable": extractable,
        "extracted_path": extracted_path,
        "not_extracted_reason": (
            not_extracted_reason if extracted_path is None else None
        ),
    }
    if package_image.resolution_warning:
        result["resolution_warning"] = package_image.resolution_warning
    if not extractable and not_extracted_reason:
        result["extraction_warning"] = not_extracted_reason
    if python_pptx_warning:
        result["inspection_warning"] = python_pptx_warning
    for attribute in ("crop_left", "crop_top", "crop_right", "crop_bottom"):
        value = getattr(shape, attribute, None)
        if value is not None:
            result[attribute] = float(value)
    return result


def shape_count(shapes: Iterable[Any]) -> int:
    count = 0
    for shape in shapes:
        count += 1
        if bool(getattr(shape, "shape_type", None)) and enum_name(shape.shape_type) == "GROUP":
            count += shape_count(shape.shapes)
    return count


def shape_dict(
    shape: Any,
    *,
    z_order: int,
    path: str,
    slide_number: int,
    slide_width: int,
    slide_height: int,
    parent_hidden: bool,
    parent_transform: AffineTransform | None,
    args: argparse.Namespace,
    limits: Limits,
    image_collector: ImageCollector,
    picture_resolver: PackagePictureResolver,
    source_part_name: str,
) -> tuple[dict[str, Any], int, int]:
    own_hidden = shape_is_hidden(shape)
    hidden = parent_hidden or own_hidden
    corners = shape_corners(shape, parent_transform)
    bounds = None if corners is None else bounds_from_corners(corners)
    position_resolved = corners is not None
    fully_off, partially_off = off_slide_state(corners, slide_width, slide_height)
    exclusion_reason: str | None = None
    if hidden and not args.include_hidden_shapes:
        exclusion_reason = "hidden_shape"
    elif fully_off and not args.include_off_slide_shapes:
        exclusion_reason = "fully_off_slide"
    elif not position_resolved and not args.include_off_slide_shapes:
        exclusion_reason = "unresolved_position"

    minimal: dict[str, Any] = {
        "shape_id": int(shape.shape_id),
        "path": path,
        "z_order": z_order,
        "hidden": hidden,
        "fully_off_slide": fully_off,
        "partially_off_slide": partially_off,
        "position_resolved": position_resolved,
        "included": exclusion_reason is None,
        "exclusion_reason": exclusion_reason,
    }
    nested_count = shape_count(shape.shapes) if enum_name(shape.shape_type) == "GROUP" else 0
    if exclusion_reason is not None:
        return minimal, 0, 1 + nested_count

    warnings: list[str] = []
    result = dict(minimal)
    result.update(
        {
            "name": add_text(shape.name, limits),
            "shape_type": enum_name(shape.shape_type),
            "position": None
            if bounds is None
            else position_dict(
                bounds,
                slide_width,
                slide_height,
                getattr(shape, "rotation", None),
            ),
            "warnings": warnings,
        }
    )
    if bool(getattr(shape, "is_placeholder", False)):
        try:
            result["placeholder_type"] = enum_name(shape.placeholder_format.type)
        except Exception as exc:
            warnings.append(f"プレースホルダー種類を取得できません: {type(exc).__name__}")
    alt_text = shape_alt_text(shape, limits)
    if alt_text:
        result["alt_text"] = alt_text
    click_action = click_action_dict(shape, limits)
    if click_action:
        result["click_action"] = click_action

    try:
        if bool(getattr(shape, "has_text_frame", False)):
            result.update(text_frame_dict(shape.text_frame, limits))
    except InspectionError:
        raise
    except Exception as exc:
        warnings.append(f"図形テキストを取得できません: {type(exc).__name__}")

    try:
        if bool(getattr(shape, "has_table", False)):
            result["table"] = table_dict(shape.table, limits)
    except InspectionError:
        raise
    except Exception as exc:
        warnings.append(f"表を取得できません: {type(exc).__name__}")

    try:
        if bool(getattr(shape, "has_chart", False)):
            result["chart"] = chart_dict(shape.chart, limits)
    except InspectionError:
        raise
    except Exception as exc:
        warnings.append(f"グラフを取得できません: {type(exc).__name__}")

    image = picture_dict(
        shape,
        limits,
        image_collector,
        picture_resolver,
        source_part_name,
        slide_number,
    )
    if image:
        result["image"] = image

    included_count = 1
    excluded_count = 0
    if enum_name(shape.shape_type) == "GROUP":
        children: list[dict[str, Any]] = []
        child_parent_transform = group_child_transform(shape, parent_transform)
        for child_index, child in enumerate(shape.shapes):
            child_path = f"{path}.{child_index + 1}"
            child_data, child_included, child_excluded = shape_dict(
                child,
                z_order=child_index,
                path=child_path,
                slide_number=slide_number,
                slide_width=slide_width,
                slide_height=slide_height,
                parent_hidden=hidden,
                parent_transform=child_parent_transform,
                args=args,
                limits=limits,
                image_collector=image_collector,
                picture_resolver=picture_resolver,
                source_part_name=source_part_name,
            )
            children.append(child_data)
            included_count += child_included
            excluded_count += child_excluded
        result["children"] = children

    if partially_off:
        warnings.append("図形の一部がスライド領域外にあります。")
    if result.get("shape_type") in {
        "DIAGRAM",
        "IGX_GRAPHIC",
        "EMBEDDED_OLE_OBJECT",
        "LINKED_OLE_OBJECT",
        "OLE_CONTROL_OBJECT",
        "FORM_CONTROL",
        "MEDIA",
        "WEB_VIDEO",
        "INK",
        "INK_COMMENT",
        "SCRIPT_ANCHOR",
    }:
        warnings.append("この図形種類は存在検出のみで、内容を完全には解析しません。")
    return result, included_count, excluded_count


def slide_title_from_shapes(shapes: list[dict[str, Any]]) -> dict[str, Any] | None:
    for shape in shapes:
        if not shape.get("included"):
            continue
        if shape.get("placeholder_type") in {"TITLE", "CENTER_TITLE"} and shape.get("text"):
            return {"text": shape["text"], "shape_id": shape["shape_id"], "path": shape["path"]}
    return None


def notes_dict(slide: Any, limits: Limits) -> dict[str, Any] | None:
    if not slide.has_notes_slide:
        return None
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    if text_frame is None:
        return {"text": "", "paragraphs": []}
    return text_frame_dict(text_frame, limits)


def document_properties_dict(properties: Any, limits: Limits) -> dict[str, Any]:
    names = (
        "title",
        "subject",
        "author",
        "keywords",
        "comments",
        "category",
        "last_modified_by",
        "created",
        "modified",
        "revision",
        "identifier",
        "language",
        "version",
    )
    result: dict[str, Any] = {}
    for name in names:
        try:
            value = getattr(properties, name)
        except Exception:
            continue
        if value in (None, ""):
            continue
        if isinstance(value, (datetime, date, time, int, float, bool)):
            result[name] = scalar_json(value)
        else:
            result[name] = add_text(value, limits)
    return result


def inventory_warnings(archive: dict[str, Any]) -> list[str]:
    warnings = [
        "図形のZオーダーは読み順ではありません。複数列、図解、重なりは視覚確認してください。",
        "スライドマスターまたはレイアウトから継承された要素は図形一覧に現れない場合があります。",
    ]
    if archive["external_relationship_count"]:
        warnings.append("外部Relationshipを検出しました。リンク先へはアクセスしていません。")
    if archive["has_vba"]:
        warnings.append("VBAを検出しました。マクロは実行、抽出、保存していません。")
    if archive["embedding_parts"]:
        warnings.append("埋め込み部品を検出しました。内容を展開または実行していません。")
    if archive["active_x_parts"]:
        warnings.append("ActiveX部品を検出しました。内容を実行していません。")
    if archive["diagram_parts"]:
        warnings.append("SmartArtまたはDiagram部品を検出しました。意味構造は完全には解析しません。")
    if archive["comment_parts"]:
        warnings.append("コメント部品を検出しました。コメント本文は抽出していません。")
    if archive["audio_video_parts"]:
        warnings.append("音声または動画を検出しました。再生、抽出、文字起こしをしていません。")
    if archive["custom_xml_parts"]:
        warnings.append("カスタムXMLを検出しました。内容は抽出していません。")
    if archive["slides_with_timing"]:
        warnings.append("アニメーションまたはタイミング設定を検出しました。表示順は解析していません。")
    if archive["slides_with_transition"]:
        warnings.append("画面切り替え設定を検出しました。見た目とタイミングは解析していません。")
    return warnings


def inspect_presentation(
    path: Path,
    source_bytes: bytes,
    args: argparse.Namespace,
) -> tuple[dict[str, Any], ImageCollector]:
    archive = inspect_archive(
        source_bytes,
        args.max_uncompressed_mb,
        macro_capable_suffix=path.suffix.casefold() == ".pptm",
    )
    try:
        presentation = Presentation(io.BytesIO(source_bytes))
    except Exception as exc:
        raise InspectionError(
            f"PowerPointのOpen XML構造を読み取れません: {type(exc).__name__}: {exc}"
        ) from exc

    slide_count_value = len(presentation.slides)
    if slide_count_value > args.max_slides:
        raise InspectionError(
            f"スライド数が上限{args.max_slides}を超えています。"
            "安全性を確認してから上限を変更してください。"
        )
    selected = sorted(set(args.slide))
    invalid_selected = [number for number in selected if number > slide_count_value]
    if invalid_selected:
        raise InspectionError(
            "存在しないスライド番号が指定されています: "
            + ", ".join(map(str, invalid_selected))
        )

    limits = Limits(
        max_slides=args.max_slides,
        max_shapes=args.max_shapes,
        max_table_cells=args.max_table_cells,
        max_chart_points=args.max_chart_points,
        max_text_chars=args.max_text_chars,
    )
    image_collector = ImageCollector(
        args.extract_images,
        args.max_extracted_images_mb * 1024 * 1024,
    )
    picture_resolver = PackagePictureResolver(source_bytes)
    slide_width = int(presentation.slide_width)
    slide_height = int(presentation.slide_height)
    slides: list[dict[str, Any]] = []
    hidden_count = 0
    included_count = 0

    for number, slide in enumerate(presentation.slides, start=1):
        hidden = slide_is_hidden(slide)
        hidden_count += int(hidden)
        explicitly_selected = number in selected
        if selected:
            included = explicitly_selected
            exclusion_reason = None if included else "not_selected"
        elif hidden and not args.include_hidden_slides:
            included = False
            exclusion_reason = "hidden_slide"
        else:
            included = True
            exclusion_reason = None

        slide_entry: dict[str, Any] = {
            "number": number,
            "slide_id": int(slide.slide_id),
            "hidden": hidden,
            "included": included,
            "exclusion_reason": exclusion_reason,
            "has_notes": bool(slide.has_notes_slide),
            "notes_included": bool(included and args.include_notes and slide.has_notes_slide),
            "layout_name": None,
            "title": None,
            "shape_count": 0,
            "included_shape_count": 0,
            "excluded_shapes": 0,
            "shapes": [],
            "notes": None,
            "warnings": [],
        }
        if not included:
            slides.append(slide_entry)
            continue

        included_count += 1
        if hidden and explicitly_selected:
            slide_entry["warnings"].append(
                "非表示スライドを--slideで明示選択したため内容を抽出しました。"
            )
        try:
            slide_entry["layout_name"] = add_text(slide.slide_layout.name, limits)
        except InspectionError:
            raise
        except Exception as exc:
            slide_entry["warnings"].append(
                f"レイアウト名を取得できません: {type(exc).__name__}"
            )

        total_shapes = shape_count(slide.shapes)
        limits.add_shapes(total_shapes)
        slide_entry["shape_count"] = total_shapes
        shape_entries: list[dict[str, Any]] = []
        source_part_name = str(slide.part.partname).lstrip("/")
        for z_order, shape in enumerate(slide.shapes):
            shape_data, shape_included, shape_excluded = shape_dict(
                shape,
                z_order=z_order,
                path=str(z_order + 1),
                slide_number=number,
                slide_width=slide_width,
                slide_height=slide_height,
                parent_hidden=False,
                parent_transform=AffineTransform(),
                args=args,
                limits=limits,
                image_collector=image_collector,
                picture_resolver=picture_resolver,
                source_part_name=source_part_name,
            )
            shape_entries.append(shape_data)
            slide_entry["included_shape_count"] += shape_included
            slide_entry["excluded_shapes"] += shape_excluded
        slide_entry["shapes"] = shape_entries
        slide_entry["title"] = slide_title_from_shapes(shape_entries)

        if args.include_notes and slide.has_notes_slide:
            slide_entry["notes"] = notes_dict(slide, limits)
        elif slide.has_notes_slide:
            slide_entry["warnings"].append(
                "発表者ノートが存在しますが、--include-notes未指定のため内容を除外しました。"
            )
        if any(shape.get("partially_off_slide") for shape in shape_entries):
            slide_entry["warnings"].append(
                "一部がスライド領域外にある図形があります。視覚確認してください。"
            )
        slides.append(slide_entry)

    document_properties = None
    if args.include_document_properties:
        document_properties = document_properties_dict(presentation.core_properties, limits)

    warnings = inventory_warnings(archive)
    if hidden_count and not args.include_hidden_slides and not selected:
        warnings.append("非表示スライドは既定で内容を除外しました。")
    if any(slide["has_notes"] for slide in slides) and not args.include_notes:
        warnings.append("発表者ノートは既定で内容を除外しました。")
    if not args.include_document_properties:
        warnings.append("文書プロパティは既定で内容を除外しました。")

    result = {
        "schema_version": SCHEMA_VERSION,
        "source": {
            "filename": escape_untrusted_text(path.name),
            "extension": path.suffix.casefold(),
            "bytes": len(source_bytes),
            "sha256": hashlib.sha256(source_bytes).hexdigest(),
        },
        "presentation": {
            "slide_count": slide_count_value,
            "included_slide_count": included_count,
            "hidden_slide_count": hidden_count,
            "slide_size": {
                "width": length_dict(slide_width),
                "height": length_dict(slide_height),
            },
            "selected_slides": selected,
            "document_properties_included": bool(args.include_document_properties),
            "document_properties": document_properties,
            "archive": archive,
            "limits": limits.as_dict(),
            "observed_counts": {
                "shapes": limits.shape_count,
                "table_cells": limits.table_cell_count,
                "chart_points": limits.chart_point_count,
                "text_characters": limits.text_char_count,
                "extracted_image_bytes": image_collector.total_bytes,
            },
        },
        "slides": slides,
        "warnings": warnings,
    }
    picture_resolver.close()
    return result, image_collector


def write_json(result: dict[str, Any], output: str | None, *, force: bool = False) -> None:
    text = json.dumps(result, ensure_ascii=False, indent=2, sort_keys=True) + "\n"
    if output is None:
        sys.stdout.write(text)
        return
    target = Path(output)
    target.parent.mkdir(parents=True, exist_ok=True)
    if not force and os.path.lexists(target):
        raise InspectionError(
            "JSONの出力先が既に存在します。上書きする場合だけ--forceを指定してください。"
        )
    file_descriptor, temporary_name = tempfile.mkstemp(
        prefix=f".{target.name}.", suffix=".tmp", dir=target.parent
    )
    temporary_path = Path(temporary_name)
    try:
        with os.fdopen(file_descriptor, "w", encoding="utf-8", newline="\n") as stream:
            stream.write(text)
            stream.flush()
            os.fsync(stream.fileno())
        if force:
            os.replace(temporary_path, target)
        else:
            try:
                os.link(temporary_path, target)
            except FileExistsError as exc:
                raise InspectionError(
                    "JSONの出力先が既に存在します。"
                    "上書きする場合だけ--forceを指定してください。"
                ) from exc
            temporary_path.unlink()
    except Exception:
        try:
            temporary_path.unlink()
        except OSError:
            pass
        raise


def main(argv: list[str] | None = None) -> int:
    args = parse_args(sys.argv[1:] if argv is None else argv)
    try:
        path = validate_args(args)
        source_bytes = path.read_bytes()
        result, image_collector = inspect_presentation(path, source_bytes, args)
        image_collector.ensure_output_does_not_collide(args.output)
        created_images = image_collector.write_all()
        try:
            write_json(result, args.output, force=args.force)
        except Exception:
            image_collector.cleanup_created(created_images)
            raise
        return 0
    except InspectionError as exc:
        print(f"エラー: {exc}", file=sys.stderr)
        return 2
    except (OSError, PermissionError) as exc:
        print(f"エラー: ファイル処理に失敗しました: {exc}", file=sys.stderr)
        return 2
    except Exception as exc:  # pragma: no cover - 予期しないライブラリエラー
        print(
            f"エラー: PowerPointの解析に失敗しました: {type(exc).__name__}: {exc}",
            file=sys.stderr,
        )
        return 4


if __name__ == "__main__":
    raise SystemExit(main())
